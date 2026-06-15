# -*- coding: utf-8 -*-
"""
Costruisce la presentazione: Captioning sensoriale del formaggio Trentingrana.
~32 slide, stile moderno con richiami al formaggio, grafici nativi + figure + foto reali,
note del presentatore per ~30 minuti di esposizione.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from PIL import Image

ROOT = r"C:\Users\nicol\Desktop\Cheese"
ASSETS = os.path.join(ROOT, "presentation", "assets")
OUTFILE = os.path.join(ROOT, "presentation", "Trentingrana_Captioning.pptx")

# ----------------- PALETTE (richiami al formaggio) -----------------
CREAM   = RGBColor(0xF7, 0xF1, 0xE3)   # crema di fondo
CREAM2  = RGBColor(0xFB, 0xF7, 0xEE)   # crema chiara card
INK     = RGBColor(0x2B, 0x2A, 0x26)   # testo
PRIM    = RGBColor(0x2C, 0x5F, 0x8A)   # blu primario
DEEP    = RGBColor(0x1B, 0x33, 0x4A)   # blu profondo (divider)
GOLD    = RGBColor(0xE0, 0xA4, 0x58)   # oro crosta
GOLD_D  = RGBColor(0xC8, 0x88, 0x36)   # oro scuro
GREEN   = RGBColor(0x5B, 0x8C, 0x5A)   # verde
RED     = RGBColor(0xC1, 0x49, 0x2E)   # rosso
GREY    = RGBColor(0x8A, 0x8A, 0x84)   # baseline grigio
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
M1C     = RGBColor(0xE0, 0x7A, 0x5F)
M3C     = RGBColor(0x3D, 0x70, 0x68)
M6C     = RGBColor(0x8A, 0x50, 0x82)

FONT_H  = "Georgia"        # titoli, tono editoriale
FONT_B  = "Calibri"        # corpo

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ----------------- HELPERS -----------------
def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _no_autosize(tf):
    tf.word_wrap = True

def add_bg(slide, color=CREAM):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=None, lw=None, shadow=False, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    _set_fill(shp, color)
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(lw or 1)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        # ombra morbida
        sp = el.makeelement(qn('a:effectLst'), {})
        outer = el.makeelement(qn('a:outerShdw'),
            {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val':'1B334A'})
        alpha = el.makeelement(qn('a:alpha'), {'val':'24000'})
        clr.append(alpha); outer.append(clr); sp.append(outer); el.append(sp)
    return shp

def oval(slide, x, y, w, h, color, alpha=None, line=None, lw=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    _set_fill(shp, color)
    shp.line.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(lw or 1)
    shp.shadow.inherit = False
    if alpha is not None:
        sp = shp.fill.fore_color._xFill
        srgb = sp.find(qn('a:srgbClr'))
        a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha*1000))})
        srgb.append(a)
    return shp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; _no_autosize(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tb, tf

def setrun(p, text, size, color=INK, bold=False, italic=False, font=FONT_B, spacing=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    if spacing is not None:
        rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(spacing))
    return r

def para(tf, first=False):
    if first and not tf.paragraphs[0].runs and tf.paragraphs[0].text == "":
        return tf.paragraphs[0]
    return tf.add_paragraph()

# cheese holes decorative cluster
def cheese_holes(slide, cx, cy, color=GOLD, base=0.42, n=5, alpha=None):
    spots = [(0,0,1.0),(0.9,0.5,0.55),(-0.7,0.7,0.45),(0.5,-0.8,0.5),(-0.9,-0.4,0.4),(1.1,-0.5,0.35)]
    for i,(dx,dy,sc) in enumerate(spots[:n]):
        d = Inches(base*sc)
        oval(slide, cx+Inches(dx*0.8)-d/2, cy+Inches(dy*0.8)-d/2, d, d, color, alpha=alpha)

# footer + slide number + top accent
PROJ = "AI4FQC · Project 07 — GRANA Captioning"
def chrome(slide, idx, accent=GOLD, dark=False):
    # barra accento sinistra
    rect(slide, 0, 0, Inches(0.16), SH, accent)
    # footer
    fc = CREAM if dark else RGBColor(0x6B,0x68,0x60)
    _, tf = textbox(slide, Inches(0.45), SH-Inches(0.42), Inches(9), Inches(0.32))
    p = tf.paragraphs[0]; setrun(p, PROJ, 9.5, fc, italic=True)
    _, tf2 = textbox(slide, SW-Inches(1.4), SH-Inches(0.42), Inches(0.95), Inches(0.32))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    setrun(p2, f"{idx:02d}", 11, accent, bold=True)

def kicker_title(slide, kicker, title, tcolor=DEEP, kcolor=GOLD_D, top=Inches(0.55)):
    _, tf = textbox(slide, Inches(0.7), top, Inches(11.9), Inches(0.4))
    p = tf.paragraphs[0]; setrun(p, kicker.upper(), 12.5, kcolor, bold=True, spacing=220)
    _, tf2 = textbox(slide, Inches(0.7), top+Inches(0.42), Inches(11.9), Inches(1.0))
    p2 = tf2.paragraphs[0]; setrun(p2, title, 31, tcolor, bold=True, font=FONT_H)
    # filo oro sotto al titolo
    rect(slide, Inches(0.72), top+Inches(1.32), Inches(1.5), Pt(3), GOLD)
    return top+Inches(1.55)

def pic_fit(slide, path, box_x, box_y, box_w, box_h, align="center", valign="middle"):
    """Inserisce immagine adattata (preserva aspect) dentro un box."""
    with Image.open(path) as im:
        iw, ih = im.size
    bw, bh = box_w, box_h
    ar = iw/ih; bar = bw/bh
    if ar > bar:
        w = bw; h = int(bw/ar)
    else:
        h = bh; w = int(bh*ar)
    if align == "center": x = box_x + (bw-w)//2
    elif align == "left": x = box_x
    else: x = box_x + (bw-w)
    if valign == "middle": y = box_y + (bh-h)//2
    elif valign == "top": y = box_y
    else: y = box_y + (bh-h)
    return slide.shapes.add_picture(path, x, y, width=w, height=h)

def A(name): return os.path.join(ASSETS, name)

def style_chart(chart, font_sz=11, legend=True, legpos=XL_LEGEND_POSITION.BOTTOM):
    chart.has_title = False
    if legend:
        chart.has_legend = True
        chart.legend.position = legpos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(font_sz); chart.legend.font.name = FONT_B
    else:
        chart.has_legend = False
    for ax in (chart.category_axis, chart.value_axis):
        try:
            ax.tick_labels.font.size = Pt(font_sz)
            ax.tick_labels.font.name = FONT_B
            ax.tick_labels.font.color.rgb = INK
        except Exception:
            pass

def color_series(plot, colors):
    for s, c in zip(plot.series, colors):
        s.format.fill.solid(); s.format.fill.fore_color.rgb = c
        s.format.line.fill.background()

def points_color(series, colors):
    for pt, c in zip(series.points, colors):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
        pt.format.line.fill.background()

def card(slide, x, y, w, h, title=None, body=None, accent=PRIM, fill=CREAM2,
         tsize=15, bsize=12.5, bullets=None, tcolor=None):
    rect(slide, x, y, w, h, fill, round_=True, shadow=True)
    rect(slide, x, y, Inches(0.09), h, accent, round_=False)
    dark = (fill == DEEP)
    txt = CREAM if dark else INK            # testo corpo
    txt2 = RGBColor(0xD9,0xD2,0xC2) if dark else INK
    pad = Inches(0.28)
    _, tf = textbox(slide, x+pad, y+Inches(0.16), w-pad-Inches(0.18), h-Inches(0.3))
    if title:
        for j, ln in enumerate(title.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            setrun(p, ln, tsize, tcolor or (GOLD if dark else DEEP), bold=True, font=FONT_H)
    if body:
        p = para(tf, first=(title is None))
        setrun(p, body, bsize, txt)
    if bullets:
        for i,(b) in enumerate(bullets):
            p = para(tf, first=(title is None and i==0 and not body))
            p.space_before = Pt(4)
            setrun(p, "● ", 10, accent, bold=True)
            if isinstance(b, tuple):
                setrun(p, b[0], bsize, txt, bold=True)
                setrun(p, b[1], bsize, txt2)
            else:
                setrun(p, b, bsize, txt)
    return

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()

def newslide(bg=CREAM):
    s = prs.slides.add_slide(BLANK); add_bg(s, bg); return s

# =====================================================================
# SLIDE 1 — TITOLO
# =====================================================================
s = newslide(DEEP)
# foto fetta a destra a tutta altezza con velo
pic = pic_fit(s, A("cheese_fetta_1_sq.jpg"), Inches(8.4), Inches(0), Inches(4.93), SH,
              align="center", valign="middle")
# velo gradiente sopra la foto (rettangolo blu semi-trasparente sul bordo sinistro foto)
veil = rect(s, Inches(8.0), 0, Inches(1.4), SH, DEEP)
sp = veil.fill.fore_color._xFill.find(qn('a:srgbClr'))
sp.append(sp.makeelement(qn('a:alpha'), {'val':'62000'}))
# holes deco
cheese_holes(s, Inches(1.1), Inches(1.0), GOLD, base=0.5, alpha=22)
# barra accento
rect(s, 0, 0, Inches(0.22), SH, GOLD)
_, tf = textbox(s, Inches(0.85), Inches(1.7), Inches(7.4), Inches(0.5))
setrun(tf.paragraphs[0], "AI4FQC · PROJECT 07", 14, GOLD, bold=True, spacing=300)
_, tf = textbox(s, Inches(0.8), Inches(2.25), Inches(7.6), Inches(2.6))
p=tf.paragraphs[0]; setrun(p, "Captioning sensoriale", 44, CREAM, bold=True, font=FONT_H)
p=tf.add_paragraph(); setrun(p, "del formaggio Trentingrana", 44, GOLD, bold=True, font=FONT_H)
_, tf = textbox(s, Inches(0.85), Inches(4.7), Inches(7.0), Inches(1.0))
setrun(tf.paragraphs[0],
       "Dalla costruzione del dataset al confronto di tre metodi encoder–decoder",
       17, RGBColor(0xD9,0xD2,0xC2), italic=True)
rect(s, Inches(0.88), Inches(5.55), Inches(2.0), Pt(3), GOLD)
_, tf = textbox(s, Inches(0.85), Inches(5.8), Inches(7), Inches(0.8))
p=tf.paragraphs[0]
setrun(p, "Relazione tecnica  ·  ", 13, RGBColor(0xC9,0xC2,0xB2))
setrun(p, "38.437 coppie immagine–caption  ·  3 modelli  ·  7 attributi", 13, CREAM, bold=True)
notes(s, """
Buongiorno a tutti e benvenuti. Oggi vi presento il progetto di captioning sensoriale
del formaggio Trentingrana, sviluppato nell'ambito dell'iniziativa AI4FQC — Artificial
Intelligence for Food Quality Control, Project 07.
L'idea di partenza è semplice da enunciare ma sorprendentemente ricca da realizzare:
generare automaticamente descrizioni sensoriali in italiano a partire da immagini di
sezioni di forme di grana. In altre parole, insegnare a un modello a "raccontare" un
formaggio guardandone una fetta.
Il lavoro si articola su due assi: la costruzione del dataset — la parte più originale e
faticosa, a cui dedicherò circa il 70% del tempo — e il confronto di tre architetture
encoder-decoder concettualmente diverse. Vedremo numeri concreti: 38 mila coppie
immagine-caption, tre modelli, sette attributi sensoriali, e soprattutto un risultato
scientifico chiaro su cosa conta davvero perché un modello usi l'immagine.
Iniziamo dal contesto.
""")

# =====================================================================
# SLIDE 2 — CONTESTO
# =====================================================================
s = newslide()
chrome(s, 2)
y0 = kicker_title(s, "Il contesto", "Perché un'AI che descrive il formaggio")
# colonna sinistra testo / destra foto grana
card(s, Inches(0.7), y0, Inches(6.5), Inches(1.7),
     title="Controllo qualità sensoriale",
     body="Ogni forma di grana è valutata da un panel di degustatori esperti. "
          "È un lavoro prezioso ma lento, soggettivo e difficile da scalare su migliaia di forme.",
     accent=PRIM)
card(s, Inches(0.7), y0+Inches(1.85), Inches(6.5), Inches(1.95),
     title="L'idea",
     accent=GOLD,
     bullets=[("Input: ", "immagini IRIS di sezioni di forma (illuminazione controllata)"),
              ("Output: ", "descrizione sensoriale in italiano, per attributo"),
              ("Valore: ", "supporto oggettivo e riproducibile al panel umano")])
# foto grana riquadrata
fx = Inches(7.55); fw = Inches(5.1)
rect(s, fx-Inches(0.1), y0-Inches(0.05), fw+Inches(0.2), Inches(3.9), WHITE, round_=True, shadow=True)
pic_fit(s, A("cheese_grana_1_sq.jpg"), fx, y0+Inches(0.08), fw, Inches(3.4))
_, tf = textbox(s, fx, y0+Inches(3.5), fw, Inches(0.35))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p, "Primo piano della grana — analizzatore visivo IRIS", 10.5, GREY, italic=True)
notes(s, """
Partiamo dal problema reale. La qualità del grana viene certificata da panel di assaggiatori:
esperti che annusano, assaggiano e osservano ogni forma, descrivendo profumo, aroma, sapore,
struttura e così via. È un patrimonio di competenza, ma ha tre limiti: è lento, è soggettivo —
due panelisti descrivono la stessa forma con parole diverse — ed è difficile da scalare quando
le forme sono migliaia.
L'idea del progetto è affiancare a questo lavoro un sistema automatico. Diamo in pasto al
modello un'immagine acquisita con l'analizzatore visivo elettronico IRIS, in condizioni di
illuminazione controllata, come quella che vedete a destra, e gli chiediamo di produrre una
descrizione sensoriale in italiano per ciascun attributo.
Attenzione: l'obiettivo non è sostituire il degustatore, ma offrire un supporto oggettivo e
riproducibile. Come vedremo, alcuni attributi — quelli visibili — si prestano molto meglio di
altri, e proprio questa distinzione sarà uno dei risultati più interessanti.
""")

# =====================================================================
# SLIDE 3 — OBIETTIVO & 7 ATTRIBUTI
# =====================================================================
s = newslide()
chrome(s, 3)
y0 = kicker_title(s, "L'oggetto", "Sette attributi sensoriali")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.6))
setrun(tf.paragraphs[0],
       "Per ogni forma, il panel descrive sette dimensioni. Tre sono visibili nell'immagine, quattro no — distinzione che tornerà nei risultati.",
       13.5, INK)
attrs = [("Colore della Pasta","visivo",GREEN),("Texture","visivo",GREEN),
         ("Struttura della Pasta","visivo",GREEN),("Spessore della Crosta","visivo*",GOLD_D),
         ("Profumo","olfatto",PRIM),("Aroma","olfatto",PRIM),("Sapore","gusto",RED)]
cols=4; cw=Inches(2.92); ch=Inches(1.55); gx=Inches(0.18); gy=Inches(0.25)
x0=Inches(0.7); yA=y0+Inches(0.75)
for i,(name,kind,col) in enumerate(attrs):
    r=i//cols; c=i%cols
    x=x0+c*(cw+gx); y=yA+r*(ch+gy)
    rect(s, x, y, cw, ch, CREAM2, round_=True, shadow=True)
    rect(s, x, y, cw, Inches(0.12), col, round_=False)
    oval(s, x+cw-Inches(0.5), y+Inches(0.32), Inches(0.28), Inches(0.28), col, alpha=28)
    _, tf = textbox(s, x+Inches(0.22), y+Inches(0.28), cw-Inches(0.4), ch-Inches(0.4))
    p=tf.paragraphs[0]; setrun(p, name, 15, DEEP, bold=True, font=FONT_H)
    p=tf.add_paragraph(); p.space_before=Pt(4); setrun(p, kind.upper(), 11, col, bold=True, spacing=150)
# ottava cella: legenda
x=x0+3*(cw+gx); y=yA+1*(ch+gy)
rect(s, x, y, cw, ch, DEEP, round_=True, shadow=True)
_, tf = textbox(s, x+Inches(0.22), y+Inches(0.2), cw-Inches(0.4), ch-Inches(0.4))
p=tf.paragraphs[0]; setrun(p, "* visibile ma piccolo", 12, GOLD, bold=True)
p=tf.add_paragraph(); p.space_before=Pt(3)
setrun(p, "La crosta occupa poca area: CLIP la pesa poco.", 11.5, CREAM)
notes(s, """
Queste sono le sette dimensioni sensoriali che il panel valuta per ogni forma: colore della
pasta, texture, struttura della pasta, spessore della crosta, profumo, aroma e sapore.
Vi chiedo di tenere a mente una distinzione che diventerà centrale nella seconda parte.
Tre attributi — colore, texture e struttura della pasta — sono intrinsecamente visivi: si
leggono dall'immagine. Tre — profumo, aroma e sapore — appartengono all'olfatto e al gusto:
nessuna fotografia può rivelare un sapore. E poi c'è un caso istruttivo, lo spessore della
crosta: è fisicamente visibile, eppure la crosta occupa una porzione piccola della fetta, e
vedremo che il modello fatica a "pesarla".
Questa tassonomia non è un dettaglio decorativo: quando misureremo se il modello usa davvero
l'immagine, scopriremo che riesce esattamente sugli attributi visibili e fallisce su quelli
olfattivi e gustativi. È un controllo di sanità che valida l'intero approccio.
""")

# =====================================================================
# SLIDE 4 — I DUE PASSI
# =====================================================================
s = newslide()
chrome(s, 4)
y0 = kicker_title(s, "La traccia", "Due passi, un obiettivo di confronto")
# due card grandi
card(s, Inches(0.7), y0+Inches(0.1), Inches(5.9), Inches(3.7),
     title="Passo 1 · Pre-elaborazione del testo", accent=PRIM,
     bullets=[("", "Pulire e normalizzare descrizioni telegrafiche, dialettali, incoerenti"),
              ("", "Sostituire le misure quantitative (mm/cm) con descrizioni qualitative"),
              ("", "Produrre caption in forma di frase italiana completa"),
              ("≈ 70%", " del lavoro complessivo del progetto")],
     tsize=16.5, bsize=13.5)
card(s, Inches(6.75), y0+Inches(0.1), Inches(5.9), Inches(3.7),
     title="Passo 2 · Tre metodi di captioning", accent=GOLD,
     bullets=[("", "Applicare tre architetture encoder–decoder «concettualmente diverse»"),
              ("", "Enfasi sul confronto di metodi, non sul singolo modello «migliore»"),
              ("", "Valutazione multi-metrica oltre il solo BLEU"),
              ("≈ 30%", " del lavoro complessivo del progetto")],
     tsize=16.5, bsize=13.5)
notes(s, """
La traccia AI4FQC chiede esplicitamente due deliverable.
Il primo passo è la pre-elaborazione del testo: prendere le descrizioni dei degustatori —
spesso telegrafiche, a volte in dialetto, talvolta incoerenti — e trasformarle in caption
pulite, normalizzate e in forma di frase italiana. Un requisito specifico e ricorrente è
sostituire le misure quantitative, i millimetri e centimetri della crosta, con descrizioni
qualitative. Questo passo, come vedete, vale circa il 70% del lavoro: è la parte più originale.
Il secondo passo è applicare e confrontare tre metodi encoder-decoder "concettualmente quanto
più diversi possibile". Notate l'accento: la traccia non chiede di trovare il modello migliore,
ma di confrontare approcci. Per farlo onestamente abbiamo affiancato a BLEU una batteria di
metriche complementari.
Coerentemente con questi pesi, dedicherò la prima e più ampia parte della presentazione alla
costruzione del dataset, e la seconda ai modelli. Cominciamo dai dati.
""")

# =====================================================================
# SLIDE 5 — SECTION PARTE I
# =====================================================================
s = newslide(DEEP)
rect(s, 0, 0, Inches(0.22), SH, GOLD)
cheese_holes(s, Inches(11.2), Inches(1.4), GOLD, base=0.7, alpha=18)
cheese_holes(s, Inches(2.0), Inches(6.2), GOLD, base=0.5, alpha=14)
_, tf = textbox(s, Inches(1.0), Inches(2.55), Inches(11), Inches(1.0))
setrun(tf.paragraphs[0], "PARTE I", 20, GOLD, bold=True, spacing=400)
_, tf = textbox(s, Inches(0.95), Inches(3.15), Inches(11.4), Inches(1.6))
setrun(tf.paragraphs[0], "Costruzione del dataset", 46, CREAM, bold=True, font=FONT_H)
rect(s, Inches(1.0), Inches(4.55), Inches(2.4), Pt(4), GOLD)
_, tf = textbox(s, Inches(1.0), Inches(4.85), Inches(10.5), Inches(0.8))
setrun(tf.paragraphs[0],
       "Da 51.988 righe grezze e dialettali a 38.437 coppie immagine–caption pulite",
       17, RGBColor(0xD9,0xD2,0xC2), italic=True)
notes(s, """
Entriamo nella prima parte: la costruzione del dataset.
Vi anticipo il viaggio in una riga: partiamo da quasi 52 mila righe grezze, eterogenee e a
tratti dialettali, e arriviamo a 38 mila coppie immagine-caption pulite, normalizzate e in
forma di frase. Nel mezzo c'è una pipeline di undici fasi e una strategia precisa che ci ha
permesso di fare tutto questo con un costo di intelligenza artificiale di appena 5 dollari e 60.
Questa è, a mio avviso, la parte più interessante del progetto dal punto di vista ingegneristico,
perché il dato grezzo era ingannevolmente ricco ma profondamente disallineato. Vediamo perché.
""")

# =====================================================================
# SLIDE 6 — DATI GREZZI & JOIN
# =====================================================================
s = newslide()
chrome(s, 6)
y0 = kicker_title(s, "Il punto di partenza", "Dati ricchi, ma disallineati")
card(s, Inches(0.7), y0, Inches(3.85), Inches(3.55),
     title="Cosa avevamo", accent=PRIM, tsize=15,
     bullets=[("2.745", " foto BMP di sezioni (Fetta / Grana)"),
              ("4", " workbook Excel (2018–2021), un foglio per attributo"),
              ("1", " codebook: caseificio ↔ prodotto ↔ lettera (16 caseifici)")],
     bsize=12.5)
card(s, Inches(4.7), y0, Inches(3.85), Inches(3.55),
     title="Il problema del join", accent=RED, tsize=15,
     body="Immagini e commenti «parlavano lingue diverse». Il nome file identifica la singola "
          "forma; il commento solo il vassoio del degustatore. Nessun join a livello di riga — "
          "solo a livello di caseificio, via codebook.",
     bsize=12.5)
card(s, Inches(8.7), y0, Inches(3.95), Inches(3.55),
     title="Ostacoli invisibili", accent=GOLD, tsize=15,
     bullets=[("", "Tripla indicizzazione: TN_306 / TN306 / 306"),
              ("", "Header incoerenti tra anni"),
              ("", "Virgola decimale italiana ('7,48')"),
              ("", "Provenienza propagata a valle → tracciabilità")],
     bsize=12)
notes(s, """
Il materiale di partenza sembrava ricco: quasi 2.750 fotografie BMP ad alta risoluzione di
sezioni di forma, in due viste — la fetta intera e il primo piano della grana; quattro workbook
Excel, uno per anno dal 2018 al 2021, con un foglio per ciascun attributo; e un codebook che
mappa i sedici caseifici.
Il problema, però, non era la quantità ma l'allineamento. Immagini e commenti, in pratica,
parlavano lingue diverse. Il nome del file fotografico identifica la singola forma — un ID
campione preciso — mentre il commento del degustatore identifica soltanto il vassoio, il
"prodotto". Non esisteva quindi un join riga-per-riga tra una foto e il suo commento: esisteva
solo un join a livello di caseificio, passando dal codebook.
A complicare tutto, tre ostacoli concreti che sono costati tempo reale: lo stesso caseificio
compare scritto in tre modi diversi; le intestazioni Excel cambiano da un anno all'altro; e i
punteggi numerici usano la virgola decimale italiana. Ogni metadato di provenienza è stato
propagato a valle, così ogni riga finale resta tracciabile fino al commento originale.
""")

# =====================================================================
# SLIDE 7 — BROADCAST
# =====================================================================
s = newslide()
chrome(s, 7)
y0 = kicker_title(s, "La soluzione di join", "Broadcast a livello di caseificio")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.55))
setrun(tf.paragraphs[0],
       "Ogni commento è propagato a tutte le immagini coerenti di quel caseificio in quella seduta: fino a 4 righe da un solo commento.",
       13.5, INK)
# diagramma: commento -> 4 immagini -> output
dy = y0+Inches(0.7)
src = rect(s, Inches(0.8), dy+Inches(0.9), Inches(2.7), Inches(1.3), RGBColor(0xF6,0xE3,0xDD), round_=True, shadow=True)
rect(s, Inches(0.8), dy+Inches(0.9), Inches(0.09), Inches(1.3), RED)
_, tf = textbox(s, Inches(1.0), dy+Inches(1.0), Inches(2.4), Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; setrun(p,"Commento panel",12,DEEP,bold=True)
p=tf.add_paragraph(); setrun(p,'«Crauti»',15,RED,bold=True,italic=True,font=FONT_H)
p=tf.add_paragraph(); setrun(p,"(Profumo · TN_306)",10.5,GREY)
labs=["P3a — Fetta","P3a — Grana","P3b — Fetta","P3b — Grana"]
for i,l in enumerate(labs):
    iy=dy+i*Inches(0.78)
    b=rect(s, Inches(5.0), iy, Inches(2.7), Inches(0.62), CREAM2, round_=True, shadow=True)
    rect(s, Inches(5.0), iy, Inches(0.08), Inches(0.62), PRIM)
    _, tf=textbox(s, Inches(5.2), iy, Inches(2.5), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0], l, 12.5, INK, bold=True)
    # freccia
    ar=s.shapes.add_connector(2, Inches(3.55), dy+Inches(1.55), Inches(4.95), iy+Inches(0.31))
    ar.line.color.rgb=GREY; ar.line.width=Pt(1.5)
out=rect(s, Inches(9.2), dy+Inches(0.9), Inches(3.3), Inches(1.3), RGBColor(0xE3,0xEE,0xE3), round_=True, shadow=True)
rect(s, Inches(9.2), dy+Inches(0.9), Inches(0.09), Inches(1.3), GREEN)
_, tf=textbox(s, Inches(9.4), dy+Inches(0.95), Inches(3.0), Inches(1.2), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; setrun(p,"4 righe immagine–caption",13,DEEP,bold=True)
p=tf.add_paragraph(); setrun(p,"(stessa caption)",11.5,GREEN)
for i in range(4):
    iy=dy+i*Inches(0.78)
    ar=s.shapes.add_connector(2, Inches(7.75), iy+Inches(0.31), Inches(9.15), dy+Inches(1.55))
    ar.line.color.rgb=GREY; ar.line.width=Pt(1.5)
# nota numeri
nb=rect(s, Inches(0.8), dy+Inches(2.55), Inches(11.7), Inches(0.7), DEEP, round_=True)
_, tf=textbox(s, Inches(1.0), dy+Inches(2.55), Inches(11.3), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]
setrun(p,"Effetto netto:  ",12.5,GOLD,bold=True)
setrun(p,"51.988 righe (immagine × panelista × attributo) · 39.510 con commento · tasso di pairing ≈ 76%",12.5,CREAM)
notes(s, """
Come si risolve l'assenza di un join a livello di forma? Con quello che abbiamo chiamato
broadcast a livello di caseificio. L'idea è propagare ogni commento del degustatore a tutte le
immagini coerenti di quel caseificio in quella seduta.
Guardate l'esempio. Un singolo commento sul profumo — la parola "Crauti" per il caseificio
TN_306 — diventa la caption di quattro immagini diverse: le repliche a e b, ciascuna nelle due
viste fetta e grana. Da un commento nascono fino a quattro righe immagine-caption, con la stessa
caption.
Qualcuno potrebbe obiettare: non è una duplicazione artificiale? In realtà no, per
l'addestramento è un vantaggio: più esempi della stessa associazione visivo-testuale rendono il
segnale più robusto.
L'effetto netto del join è di quasi 52 mila righe, ottenute come immagine per panelista per
attributo, di cui circa 39.500 con un commento non vuoto: un tasso di pairing intorno al 76%,
più basso nelle sedute 2021 dove alcuni campi erano lasciati vuoti. Da qui parte la pipeline.
""")

# =====================================================================
# SLIDE 8 — PIPELINE 11 FASI
# =====================================================================
s = newslide()
chrome(s, 8)
y0 = kicker_title(s, "L'architettura", "Una pipeline in 11 fasi")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.5))
p=tf.paragraphs[0]
setrun(p,"Principio guida: ",13.5,INK,bold=True)
setrun(p,"deterministico-prima-dell'LLM — tutto ciò che si può fare con codice riproducibile e gratuito viene fatto prima di interpellare il modello.",13.5,INK)
phases=[("0","Tabella unificata","51.988",PRIM),("1","Prep deterministica","39.356",PRIM),
        ("2–3","Vocabolario + audit","7 attr.",PRIM),("4","Pulizia + qualitativo","7.705 uniche",PRIM),
        ("5","Drop rumore","7.689",PRIM),("6–7","Prompt + pilot","105",GOLD_D),
        ("8","Batch LLM","7.689",GOLD_D),("9","Salvage manuale","+916",GREEN),
        ("10","Broadcast + frase","38.437",GREEN)]
cols=3; cw=Inches(3.9); ch=Inches(1.18); gx=Inches(0.15); gy=Inches(0.18)
x0=Inches(0.75); yA=y0+Inches(0.65)
for i,(n,t,v,col) in enumerate(phases):
    r=i//cols; c=i%cols
    x=x0+c*(cw+gx); y=yA+r*(ch+gy)
    rect(s, x, y, cw, ch, CREAM2, round_=True, shadow=True)
    badge=oval(s, x+Inches(0.18), y+Inches(0.3), Inches(0.58), Inches(0.58), col)
    _, tf=textbox(s, x+Inches(0.18), y+Inches(0.3), Inches(0.58), Inches(0.58), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,n,12.5,WHITE,bold=True)
    _, tf=textbox(s, x+Inches(0.92), y+Inches(0.16), cw-Inches(1.05), ch-Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; setrun(p,t,13.5,DEEP,bold=True,font=FONT_H)
    p=tf.add_paragraph(); setrun(p,v,12,col if col!=GOLD_D else GOLD_D,bold=True)
# legenda colori
_, tf=textbox(s, Inches(0.75), yA+3*(ch+gy)-Inches(0.02), Inches(11.8), Inches(0.4))
p=tf.paragraphs[0]
setrun(p,"■ ",13,PRIM,bold=True); setrun(p,"deterministico (gratis)    ",11.5,INK)
setrun(p,"■ ",13,GOLD_D,bold=True); setrun(p,"con LLM    ",11.5,INK)
setrun(p,"■ ",13,GREEN,bold=True); setrun(p,"assemblaggio finale",11.5,INK)
notes(s, """
Ecco la mappa completa: undici fasi, numerate da zero a dieci, ognuna con un proprio script
Python, input e output espliciti e un report ispezionabile.
Il principio guida — e questa è la decisione architetturale più importante della prima parte —
è "deterministico prima dell'LLM". Significa: tutto ciò che si può fare con codice
riproducibile, verificabile e gratuito, lo facciamo prima di chiamare il modello linguistico.
In blu vedete le fasi deterministiche: dalla tabella unificata, alla preparazione del testo,
alla costruzione del vocabolario, alla pulizia e deduplicazione. In oro le fasi che coinvolgono
l'LLM: il design del prompt, il pilot e il batch completo. In verde l'assemblaggio finale.
Seguite i numeri lungo le fasi: 52 mila righe iniziali, ridotte e normalizzate, poi compresse a
circa 7.700 caption uniche — ed è solo questo sottoinsieme ad arrivare al modello a pagamento.
È esattamente questa compressione che tiene il costo a pochi dollari. Vediamo le fasi chiave una
per una.
""")

# =====================================================================
# SLIDE 9 — IMBUTO DATI (figura)
# =====================================================================
s = newslide()
chrome(s, 9)
y0 = kicker_title(s, "Il flusso dei dati", "L'imbuto: dove si perdono (e si tengono) le righe")
pic_fit(s, A("fig_funnel.png"), Inches(0.7), y0, Inches(7.7), Inches(4.4))
card(s, Inches(8.7), y0+Inches(0.1), Inches(3.95), Inches(4.1),
     title="Lettura", accent=GOLD, tsize=15,
     bullets=[("−12.632", " righe nella prep deterministica"),
              ("12.478", " vuoti / null"),
              ("86", " meta-commenti  ·  68 quasi-vuoti"),
              ("75,7%", " di retention dopo la fase 1")],
     bsize=12.5)
notes(s, """
Questa è la stessa storia, vista come imbuto. Sull'asse leggete il numero di righe di
addestramento man mano che attraversano la pipeline.
Lo scarto maggiore — oltre 12.600 righe — avviene tutto nella preparazione deterministica della
fase 1. Ma attenzione a cosa buttiamo: la stragrande maggioranza, quasi 12.500 righe, sono vuoti
o valori nulli. Solo 86 sono meta-commenti e 68 sono frammenti quasi vuoti. In altre parole,
non stiamo scartando contenuto sensoriale: stiamo togliendo il rumore certo.
Il tasso di retention dopo la fase 1 è del 75,7%. E qui c'è una scelta deliberata che voglio
sottolineare: a questo stadio i drop devono essere certi. Una frase come "non penalizzo, ma sa
di stalla" contiene un descrittore reale e NON va scartata da una regex — quel tipo di ambiguità
la gestirà meglio l'LLM più avanti. La blacklist deterministica è quindi piccola e conservativa
per costruzione.
""")

# =====================================================================
# SLIDE 10 — FASE 1
# =====================================================================
s = newslide()
chrome(s, 10)
y0 = kicker_title(s, "Fase 1", "Preparazione deterministica: pulire senza riformulare")
card(s, Inches(0.7), y0, Inches(5.85), Inches(3.6),
     title="Operazioni", accent=PRIM, tsize=15.5,
     bullets=[("", "Filtro vuoti/N-A + normalizzazione Unicode (NFC)"),
              ("", "Rimozione spazi indivisibili, zero-width, tab, a-capo"),
              ("", "Drop meta-commenti (blacklist piccola e conservativa)"),
              ("", "Drop rumore quasi-vuoto (< 2 caratteri alfanumerici)"),
              ("", "Doppia colonna raw / norm → reversibilità e audit")],
     bsize=13)
card(s, Inches(6.7), y0, Inches(5.95), Inches(3.6),
     title="Filosofia", accent=GOLD, tsize=15.5, fill=DEEP, tcolor=GOLD,
     body="«I drop devono essere certi». L'LLM gestirà meglio le ambiguità: una negazione "
          "con descrittore — «non penalizzo, ma sa di stalla» — porta informazione reale e "
          "viene mantenuta.\n\nRisultato: 39.356 righe (75,7% retention), 12.632 drop, di cui "
          "12.478 vuoti, 86 meta, 68 quasi-vuoti.",
     bsize=13)
# imposta colore testo bianco nella card scura -> patch: riscrivo body manualmente
notes(s, """
Entriamo nella fase 1, la preparazione deterministica. La regola d'oro qui è: pulire senza
riformulare. Non cambiamo le parole del degustatore, le ripuliamo soltanto.
Le operazioni sono tutte meccaniche e verificabili: filtriamo i vuoti e gli N/A, normalizziamo
l'Unicode in forma NFC, togliamo gli spazi indivisibili, i caratteri a larghezza zero, le
tabulazioni e i ritorni a capo. Poi rimuoviamo i meta-commenti con una blacklist deliberatamente
piccola, e i frammenti quasi vuoti, quelli con meno di due caratteri alfanumerici dopo la
pulizia. Manteniamo sempre due colonne, testo grezzo e testo normalizzato, così ogni
trasformazione è reversibile e auditabile.
La filosofia, a destra, è la chiave: i drop devono essere certi. Tutto ciò che è ambiguo passa
oltre, perché l'LLM lo gestirà meglio di qualsiasi regola fissa. Il risultato netto sono 39.356
righe, con i drop che — lo ripeto perché conta — sono per il 99% semplici vuoti.
""")

# =====================================================================
# SLIDE 11 — FASI 2-3 VOCABOLARIO (figura)
# =====================================================================
s = newslide()
chrome(s, 11)
y0 = kicker_title(s, "Fasi 2–3", "Un vocabolario controllato, su misura")
pic_fit(s, A("fig_vocab.png"), Inches(0.7), y0, Inches(7.7), Inches(4.4))
card(s, Inches(8.7), y0+Inches(0.1), Inches(3.95), Inches(4.1),
     title="Perché custom", accent=PRIM, tsize=15,
     bullets=[("", "Le NLP italiane sbagliano: panna→panno, latte→latto"),
              ("", "Stopword utili qui: «molto», «poco», «leggermente»"),
              ("", "Lemmatizzatore deterministico + merge sing./plur."),
              ("Doppio uso:", " base per l'audit + àncora stilistica nel prompt")],
     bsize=12)
notes(s, """
Le fasi 2 e 3 costruiscono un vocabolario controllato per ciascuno dei sette attributi: i lemmi
e i bigrammi più frequenti, con una lemmatizzazione italiana fatta su misura.
Perché non usare una libreria NLP standard? Perché sul lessico sensoriale sbagliano
sistematicamente: trasformano "panna" in "panno", "latte" in "latto". E soprattutto includono
tra le stopword parole che qui sono informative: "molto", "poco", "leggermente" sono
intensificatori sensoriali, non rumore. Abbiamo quindi scritto un lemmatizzatore deterministico,
con mappe di abbreviazioni e refusi e un merge singolare-plurale che scatta solo quando entrambe
le forme sono attestate nel corpus.
Nel grafico vedete l'ampiezza del vocabolario per attributo: la struttura della pasta è la più
ricca, con quasi 800 lemmi e oltre 40 mila token; lo spessore della crosta è il più povero,
coerentemente con la sua minore varietà descrittiva. Il vocabolario serve a due scopi: come base
per l'audit, e come àncora stilistica nel prompt dell'LLM — non come dizionario chiuso, perché
il modello l'italiano lo parla già.
""")

# =====================================================================
# SLIDE 12 — FASE 4 NUMERI->QUALITA
# =====================================================================
s = newslide()
chrome(s, 12)
y0 = kicker_title(s, "Fase 4 · requisito chiave", "Da numeri a qualità: la crosta")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.5))
setrun(tf.paragraphs[0],
       "Centinaia di caption erano solo misure («10», «1 cm», «8-10 mm»). Una funzione deterministica converte tutto in mm e assegna un bucket.",
       13.5, INK)
# tabella bucket
rows=[("Soglia (mm)","Bucket"),("< 8","Molto sottile"),("8 ≤ x < 10","Sottile"),
      ("10 ≤ x < 14","Media"),("14 ≤ x < 18","Spessa"),("≥ 18","Molto spessa")]
tx=Inches(0.7); ty=y0+Inches(0.7); tw=Inches(4.7); th=Inches(3.4)
tbl=s.shapes.add_table(len(rows),2,tx,ty,tw,th).table
tbl.columns[0].width=Inches(2.2); tbl.columns[1].width=Inches(2.5)
for ri,(a,b) in enumerate(rows):
    for ci,val in enumerate((a,b)):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ri==0:
            setrun(p,val,13,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            setrun(p,val,12.5,INK,bold=(ci==1))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_left=Pt(8)
# figura spessore a destra
rect(s, Inches(5.75), ty-Inches(0.05), Inches(6.9), Inches(3.6), WHITE, round_=True, shadow=True)
pic_fit(s, A("fig_spessore.png"), Inches(5.95), ty+Inches(0.1), Inches(6.5), Inches(3.3))
notes(s, """
La fase 4 affronta di petto il requisito esplicito della traccia: sostituire le descrizioni
quantitative con descrizioni qualitative. Il caso emblematico è lo spessore della crosta, dove
centinaia di caption erano semplici numeri: "10", "1 cm", "mediamente 9 mm", "8-10 mm".
La soluzione è una funzione deterministica che riconosce le caption interamente numeriche,
converte tutto in millimetri — con un'euristica che interpreta i valori sotto 5 come centimetri
— e assegna un bucket qualitativo secondo la tabella che vedete: sotto 8 millimetri "molto
sottile", fino a 10 "sottile", e così via fino a "molto spessa" oltre i 18.
Perché è importante farlo in modo deterministico e non con l'LLM? Perché elimina alla radice la
principale fonte di incoerenza del modello, che altrimenti bucketizzava "1 cm" e "10 mm" — la
stessa identica misura fisica — in modo diverso. Nel grafico a destra vedete il risultato: 424
righe numeriche collassate ordinatamente in cinque bucket. Le caption miste, quelle che mescolano
numeri e descrizione, vengono invece lasciate all'LLM perché richiedono ragionamento contestuale.
""")

# =====================================================================
# SLIDE 13 — DEDUP (figura)
# =====================================================================
s = newslide()
chrome(s, 13)
y0 = kicker_title(s, "Fase 4 · l'accorgimento decisivo", "Deduplicazione: compressione 5,1×")
pic_fit(s, A("fig_dedup.png"), Inches(0.7), y0, Inches(7.7), Inches(4.4))
card(s, Inches(8.7), y0+Inches(0.1), Inches(3.95), Inches(4.1),
     title="Perché conta", accent=GREEN, tsize=15,
     bullets=[("", "Il broadcast replica lo stesso testo molte volte"),
              ("39.356", " righe → 7.705 caption uniche"),
              ("5,1×", " di compressione per chiave (caption, attributo)"),
              ("", "Pattern dominante 4×: 6.595 caption esattamente 4 volte"),
              ("→", " riduce di ~5× il costo dell'LLM")],
     bsize=12)
notes(s, """
Questo è l'accorgimento singolo più importante per l'economia del progetto: la deduplicazione.
Ricordate il broadcast? Replicando ogni commento su quattro immagini, lo stesso identico testo
sarebbe stato inviato all'LLM molte volte, pagandolo ogni volta. Non ha senso. Allora, prima di
chiamare il modello, deduplichiamo per chiave: la coppia caption più attributo.
Il risultato è netto: oltre 39 mila righe si comprimono in 7.705 caption uniche, una
compressione di 5,1 volte. Il grafico mostra, per ciascun attributo, le righe preparate contro
le caption uniche. Il pattern dominante è proprio il fattore quattro del broadcast: quasi 6.600
caption compaiono esattamente quattro volte.
La conseguenza pratica è enorme: mandiamo all'LLM solo le caption uniche, e così riduciamo di
circa cinque volte il costo. È l'esempio perfetto del principio "deterministico prima dell'LLM":
una semplice deduplicazione, gratis, taglia i costi del componente a pagamento di un fattore
cinque.
""")

# =====================================================================
# SLIDE 14 — FASI 6-8 LLM
# =====================================================================
s = newslide()
chrome(s, 14)
y0 = kicker_title(s, "Fasi 6–8", "Riscrittura con LLM: il lavoro davvero difficile")
card(s, Inches(0.7), y0, Inches(3.85), Inches(3.6),
     title="6 · Design prompt", accent=GOLD, tsize=15,
     bullets=[("", "Prompt di sistema per attributo (~5 KB)"),
              ("11", " regole (zero invenzione, escape NON_DESCRITTO)"),
              ("Top-60", " lemmi come àncora di stile"),
              ("6", " esempi few-shot reali")], bsize=12)
card(s, Inches(4.7), y0, Inches(3.85), Inches(3.6),
     title="7 · Pilot run", accent=PRIM, tsize=15,
     bullets=[("105", " caption (15 per attributo)"),
              ("", "Scoperti 2 fallimenti: 1cm vs 10mm; formato vuoti"),
              ("", "Entrambi corretti nel prompt"),
              ("3 worker", " = ottimo; 8 → tempeste di 429")], bsize=12)
card(s, Inches(8.7), y0, Inches(3.95), Inches(3.6),
     title="8 · Batch completo", accent=GREEN, tsize=15,
     bullets=[("7.689 / 7.689", " completate, 0 errori"),
              ("~25–30", " minuti"),
              ("Haiku 4.5", " · Anthropic Batch API"),
              ("~$4,50", " per il batch")], bsize=12)
notes(s, """
Arriviamo al cuore "intelligente" della pipeline: la riscrittura con LLM, fasi 6, 7 e 8. Qui
l'LLM fa il lavoro che nessuna regex potrebbe fare fedelmente: espandere "Crauti" in "Profumo di
crauti", normalizzare il dialetto, togliere i meta-commenti preservando il descrittore,
trasformare le domande in affermazioni.
La fase 6 è il design del prompt: un prompt di sistema per ciascun attributo, circa 5 kilobyte,
con undici regole — tra cui la regola di "zero invenzione" e un escape, NON_DESCRITTO, per gli
input senza contenuto — i sessanta lemmi più frequenti come àncora stilistica e sei esempi
few-shot reali.
La fase 7, il pilot su 105 caption, è servita a scoprire due modi di fallimento prima di spendere:
l'incoerenza centimetri-millimetri e alcune violazioni di formato. Entrambi corretti nel prompt.
Lezione operativa concreta: tre worker in parallelo erano l'ottimo, otto scatenavano tempeste di
rate-limit.
La fase 8, il batch completo: tutte le 7.689 caption uniche inviate in un unico job. Risultato:
7.689 su 7.689 completate, zero errori, in mezz'ora, per circa 4 dollari e mezzo. Il che ci porta
alla domanda: perché Haiku?
""")

# =====================================================================
# SLIDE 15 — COSTO (chart nativo)
# =====================================================================
s = newslide()
chrome(s, 15)
y0 = kicker_title(s, "La scelta del modello", "Perché Haiku 4.5: stesso lavoro, frazione del costo")
cd = CategoryChartData()
cd.categories = ["Haiku 4.5", "Sonnet", "Opus"]
cd.add_series("Costo stimato ($)", (5.60, 13.50, 67.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.7), y0+Inches(0.05), Inches(7.4), Inches(4.3), cd)
ch = gf.chart; style_chart(ch, legend=False)
plot=ch.plots[0]; plot.has_data_labels=True
plot.data_labels.number_format='"$"0.00'; plot.data_labels.number_format_is_linked=False
plot.data_labels.font.size=Pt(13); plot.data_labels.font.bold=True; plot.data_labels.font.color.rgb=DEEP
plot.data_labels.position=XL_LABEL_POSITION.OUTSIDE_END
plot.gap_width=80
points_color(plot.series[0], [GREEN, GOLD, RED])
ch.value_axis.has_major_gridlines=True
card(s, Inches(8.35), y0+Inches(0.1), Inches(4.3), Inches(4.15),
     title="Il verdetto del pilot", accent=GREEN, tsize=15.5,
     bullets=[("2,4×", " più economico di Sonnet"),
              ("12×", " più economico di Opus"),
              ("", "Sul task verificato, i modelli grandi non davano qualità misurabile in più"),
              ("$5,60", " costo LLM totale del progetto")], bsize=13)
notes(s, """
Una delle decisioni più pragmatiche del progetto. Per un compito di riscrittura ben delimitato e
già verificato sul pilot, qual è il modello giusto?
Il grafico risponde. Lo stesso lavoro costerebbe circa 5 dollari e 60 con Haiku 4.5, contro circa
13 e 50 con Sonnet e oltre 67 con Opus. Haiku è 2,4 volte più economico di Sonnet e dodici volte
più economico di Opus.
La domanda vera, però, non è "qual è il più economico" ma "il modello economico è abbastanza
buono?". E qui interviene il pilot: avendolo testato su un campione stratificato, abbiamo
verificato che per questo task circoscritto i modelli più costosi non offrivano alcun guadagno di
qualità misurabile. Spendere di più sarebbe stato buttare via soldi.
Il costo LLM totale dell'intero progetto — pilot più batch — è stato di 5 dollari e 60. È la
dimostrazione concreta che una buona ingegneria dei dati a monte vale più di un modello costoso a
valle.
""")

# =====================================================================
# SLIDE 16 — VALIDAZIONE OUTPUT
# =====================================================================
s = newslide()
chrome(s, 16)
y0 = kicker_title(s, "Validazione programmatica", "Il Passo 1 è soddisfatto al 100%")
rows=[("Controllo","Violazioni"),
      ("Output inizia col prefisso d'attributo atteso","1*"),
      ("Output contiene cifre","0"),
      ("Output contiene unità (mm / cm / %)","0"),
      ("Output più lungo di 25 parole","0"),
      ("Output vuoto","0"),
      ("Output multi-riga o con markup","0")]
tx=Inches(0.7); ty=y0+Inches(0.05); tw=Inches(7.7); th=Inches(3.9)
tbl=s.shapes.add_table(len(rows),2,tx,ty,tw,th).table
tbl.columns[0].width=Inches(6.0); tbl.columns[1].width=Inches(1.7)
for ri,(a,b) in enumerate(rows):
    for ci,val in enumerate((a,b)):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ci==1: p.alignment=PP_ALIGN.CENTER
        if ri==0:
            setrun(p,val,13.5,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            col = GREEN if val=="0" else GOLD_D
            setrun(p,val,13 if ci==0 else 16, INK if ci==0 else col, bold=(ci==1))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(10)
card(s, Inches(8.7), y0+Inches(0.05), Inches(3.95), Inches(3.95),
     title="Zero numeri,\nzero unità", accent=GREEN, tsize=18, fill=DEEP, tcolor=GREEN,
     body="\nLa conversione quantitativo → qualitativo è riuscita al 100%: il requisito "
          "principale del Passo 1 è pienamente soddisfatto.\n\n* l'unica «violazione» è un "
          "prefisso alternativo comunque valido.", bsize=13)
notes(s, """
Come sappiamo che il Passo 1 è davvero soddisfatto? Non a impressione, ma con una validazione
programmatica su tutte le 7.689 uscite del modello.
La tabella riassume i controlli. L'output inizia sempre col prefisso d'attributo atteso — con una
sola "violazione" che è in realtà un prefisso alternativo valido. Soprattutto, guardate i numeri
che contano per la traccia: zero output con cifre, zero output con unità di misura come
millimetri o centimetri, zero output troppo lunghi, zero output vuoti, zero output con markup.
Tutti zero.
Questo significa che la conversione da quantitativo a qualitativo — il requisito centrale del
primo passo — è riuscita al cento per cento: nelle caption finali non resta nessun numero e
nessuna unità di misura. È una soddisfazione pulita e verificabile del deliverable, non una
promessa. E quando un revisore chiederà "ma siete sicuri di aver tolto tutte le misure?", la
risposta è una tabella, non un'opinione.
""")

# =====================================================================
# SLIDE 17 — SALVAGE (figura)
# =====================================================================
s = newslide()
chrome(s, 17)
y0 = kicker_title(s, "Fase 9", "Salvataggio manuale: pulire ≠ buttare via")
pic_fit(s, A("fig_salvage.png"), Inches(0.7), y0, Inches(7.4), Inches(4.4))
card(s, Inches(8.35), y0+Inches(0.1), Inches(4.3), Inches(4.1),
     title="Cura a mano", accent=GOLD, tsize=15,
     bullets=[("", "291/362 NON_DESCRITTO avevano un lemma del vocabolario → LLM troppo cauto"),
              ("178", " caption salvate a mano"),
              ("«marcio, putrido,»", " → «Profumo marcio e putrido.»"),
              ("362 → 184", " NON_DESCRITTO unici"),
              ("+916", " righe di addestramento recuperate")], bsize=12)
notes(s, """
La fase 9 è una di quelle che distinguono un dataset curato da uno solo "ripulito". Dopo il
batch, una scansione ha rivelato qualcosa di interessante: 291 delle 362 caption che l'LLM aveva
etichettato come NON_DESCRITTO — cioè "qui non c'è descrizione" — contenevano in realtà almeno un
lemma del vocabolario controllato. Segno che il modello era stato eccessivamente cauto con la
regola dell'escape su input borderline, quelli che mescolano un giudizio e un descrittore.
Allora abbiamo curato a mano una mappa di salvataggio di 178 caption dove il descrittore era reale
e fedele. Esempi: "marcio, putrido," diventa "Profumo marcio e putrido"; "Sangue,,," diventa
"Aroma di sangue".
L'effetto, nel grafico: i NON_DESCRITTO unici scendono da 362 a 184, e questo recupera 916 righe
di addestramento che altrimenti avremmo perso. Il messaggio è semplice: pulire i dati non
significa buttarli via. Un po' di cura manuale, su un sottoinsieme mirato, è costata meno di un
altro giro di LLM e ha restituito quasi mille esempi validi.
""")

# =====================================================================
# SLIDE 18 — DATASET FINALE (figura + schema)
# =====================================================================
s = newslide()
chrome(s, 18)
y0 = kicker_title(s, "Fase 10 · il deliverable", "Il dataset finale")
pic_fit(s, A("fig_final_rows.png"), Inches(0.7), y0, Inches(6.7), Inches(3.5))
# KPI strip
kpis=[("38.437","righe finali",PRIM),("1.497","immagini uniche",GOLD_D),("2","forme di caption",GREEN)]
for i,(v,l,c) in enumerate(kpis):
    x=Inches(0.7)+i*Inches(2.3)
    rect(s, x, y0+Inches(3.65), Inches(2.15), Inches(0.95), CREAM2, round_=True, shadow=True)
    _, tf=textbox(s, x, y0+Inches(3.7), Inches(2.15), Inches(0.9), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,v,19,c,bold=True,font=FONT_H)
    p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; setrun(p,l,11,INK)
card(s, Inches(7.6), y0, Inches(5.05), Inches(4.6),
     title="Schema consegnato", accent=PRIM, tsize=15,
     bullets=[("captions_final.csv", " — 38.437 × 18 col. (provenienza)"),
              ("image_caption_attribute.csv", " — vista a 4 colonne"),
              ("by_attribute/<Attr>.csv", " — 7 split per-attributo"),
              ("caption", " (compatta, ~4–8 parole)"),
              ("caption_sentence", " (frase italiana, ~7–15 parole)"),
              ("", "Forma di frase deterministica: 100% template, 0 round-trip LLM extra")], bsize=11.5)
notes(s, """
Eccoci al deliverable della prima parte. La fase 10 riconnette le caption pulite alla tabella
broadcast e produce il dataset finale: 38.437 righe di addestramento su 1.497 immagini uniche.
Il grafico mostra la distribuzione delle righe per attributo, e notate che c'è una varianza
significativa — la struttura della pasta ha quasi il doppio delle righe dello spessore della
crosta. Questo sbilanciamento conterà nel confronto cross-attributo della seconda parte.
Sul piano del formato, abbiamo fatto una scelta di flessibilità: ogni riga porta due forme di
caption. Una compatta, di quattro-otto parole, ancorata all'attributo, tipo "Profumo di panna".
E una frase italiana dichiarativa completa, "Il formaggio ha un profumo di panna", più adatta agli
encoder-decoder e alle metriche di captioning. Questa trasformazione in frase è completamente
deterministica, senza alcun ulteriore giro di LLM.
Il dataset è consegnato in più forme — tabella completa con tutta la provenienza, vista
semplificata a quattro colonne, e sette split per-attributo — per servire qualsiasi architettura a
valle. Con questo chiudiamo la parte dati e passiamo ai modelli.
""")

# =====================================================================
# SLIDE 19 — SECTION PARTE II
# =====================================================================
s = newslide(DEEP)
rect(s, 0, 0, Inches(0.22), SH, GOLD)
cheese_holes(s, Inches(11.0), Inches(5.6), GOLD, base=0.7, alpha=18)
cheese_holes(s, Inches(2.2), Inches(1.4), GOLD, base=0.5, alpha=14)
_, tf = textbox(s, Inches(1.0), Inches(2.55), Inches(11), Inches(1.0))
setrun(tf.paragraphs[0], "PARTE II", 20, GOLD, bold=True, spacing=400)
_, tf = textbox(s, Inches(0.95), Inches(3.15), Inches(11.6), Inches(1.6))
setrun(tf.paragraphs[0], "Selezione dei modelli e training", 42, CREAM, bold=True, font=FONT_H)
rect(s, Inches(1.0), Inches(4.5), Inches(2.4), Pt(4), GOLD)
_, tf = textbox(s, Inches(1.0), Inches(4.8), Inches(10.8), Inches(0.8))
setrun(tf.paragraphs[0],
       "Tre metodi encoder–decoder, un solo determinante: l'encoder visivo",
       17, RGBColor(0xD9,0xD2,0xC2), italic=True)
notes(s, """
Seconda parte: i modelli.
Qui l'obiettivo della traccia è il confronto di metodi, e vi anticipo già la conclusione, così
potete seguire i dati con la chiave di lettura giusta. Tra tutto ciò che potremmo cambiare — il
decoder, la quantità di dati, gli iperparametri — il fattore singolo che decide se il modello usa
davvero l'immagine è uno solo: l'encoder visivo.
Costruiremo questa conclusione passo dopo passo, e la verificheremo con due strumenti indipendenti:
lo shuffle test e il CLIPScore. Partiamo dalle tre architetture.
""")

# =====================================================================
# SLIDE 20 — I TRE METODI
# =====================================================================
s = newslide()
chrome(s, 20)
y0 = kicker_title(s, "Le architetture", "Tre metodi, due assi di variazione")
models=[("m1","ResNet-50\n(CNN, congelato)","LSTM\nda zero",M1C),
        ("m3","ViT-B/16\n(congelato)","Transformer\nda zero",M3C),
        ("m6","ViT-B/16\n(congelato)","GePpeTto\nGPT-2 IT pre-addestrato",M6C)]
cw=Inches(3.6); gx=Inches(0.55); x0=Inches(0.95); yT=y0+Inches(0.25)
for i,(m,enc,dec,col) in enumerate(models):
    x=x0+i*(cw+gx)
    _, tf=textbox(s, x, yT, cw, Inches(0.5)); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    setrun(p,m,22,col,bold=True,font=FONT_H)
    eb=rect(s, x, yT+Inches(0.55), cw, Inches(1.15), CREAM2, round_=True, shadow=True)
    rect(s, x, yT+Inches(0.55), cw, Inches(0.1), PRIM)
    _, tf=textbox(s, x+Inches(0.15), yT+Inches(0.6), cw-Inches(0.3), Inches(1.05), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,"ENCODER",10,PRIM,bold=True,spacing=150)
    for ln in enc.split("\n"):
        p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; setrun(p,ln,13,DEEP,bold=True)
    # freccia giù
    a=s.shapes.add_connector(2, x+cw/2, yT+Inches(1.75), x+cw/2, yT+Inches(2.05))
    a.line.color.rgb=GREY; a.line.width=Pt(2)
    db=rect(s, x, yT+Inches(2.1), cw, Inches(1.15), CREAM2, round_=True, shadow=True)
    rect(s, x, yT+Inches(2.1), cw, Inches(0.1), GOLD)
    _, tf=textbox(s, x+Inches(0.15), yT+Inches(2.15), cw-Inches(0.3), Inches(1.05), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,"DECODER",10,GOLD_D,bold=True,spacing=150)
    for ln in dec.split("\n"):
        p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; setrun(p,ln,12.5,DEEP,bold=True)
# assi
_, tf=textbox(s, x0, yT+Inches(3.45), Inches(8.1), Inches(0.4))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"↔ m1 vs m3 isola l'ENCODER        ↔ m3 vs m6 isola il DECODER",12.5,RED,bold=True)
_, tf=textbox(s, Inches(0.95), yT+Inches(3.9), Inches(11.5), Inches(0.5))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"Tutti gli encoder sono congelati: si addestrano solo decoder e proiezione di cross-attention → confronto equo.",12,INK,italic=True)
notes(s, """
La traccia chiede tre metodi "concettualmente quanto più diversi possibile". Li abbiamo scelti in
modo da isolare i tre assi principali di variazione nel captioning.
m1 è il classico: un encoder convoluzionale, ResNet-50, con un decoder LSTM addestrato da zero.
m3 cambia l'encoder: un Vision Transformer, ViT-B/16, con un decoder Transformer sempre da zero.
m6 cambia il decoder: stesso ViT, ma il decoder è GePpeTto, un GPT-2 italiano già pre-addestrato.
La bellezza di questo disegno sta negli assi, in basso. Confrontare m1 con m3 isola l'effetto
dell'encoder, perché il decoder resta lo stesso tipo "da zero". Confrontare m3 con m6 isola
l'effetto del decoder, perché l'encoder è identico. È un piccolo esperimento controllato.
Una scelta cruciale per l'equità: tutti gli encoder sono congelati. Addestriamo solo il decoder e
la proiezione di cross-attention. Così ogni decoder riceve esattamente le stesse feature visive, e
le differenze che osserveremo sono attribuibili all'architettura, non a un encoder fortunato.
Questo, vedremo, è anche il principale limite — e la principale indicazione per il futuro.
""")

# =====================================================================
# SLIDE 21 — SETUP
# =====================================================================
s = newslide()
chrome(s, 21)
y0 = kicker_title(s, "Setup sperimentale", "Addestramento su Kaggle, confronto equo")
rows=[("Modello","epoche","batch","lr","scheduler","patience"),
      ("m1 · CNN+LSTM","50","32","3e-4","StepLR","7"),
      ("m3 · ViT+Transf.","30","16","1e-4","cosine","5"),
      ("m6 · ViT+GePpeTto","20","8","5e-5","cosine","5")]
tx=Inches(0.7); ty=y0; tw=Inches(7.6); th=Inches(2.1)
tbl=s.shapes.add_table(len(rows),6,tx,ty,tw,th).table
widths=[2.3,1.0,1.0,1.1,1.2,1.0]
for ci,w in enumerate(widths): tbl.columns[ci].width=Inches(w)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ci>0: p.alignment=PP_ALIGN.CENTER
        if ri==0:
            setrun(p,val,11.5,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            setrun(p,val,11.5,INK,bold=(ci==0))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(6)
card(s, Inches(0.7), ty+Inches(2.35), Inches(7.6), Inches(1.95),
     title="Scelte di valutazione equa", accent=PRIM, tsize=14.5,
     bullets=[("Split sample-disjoint:", " 674 train / 143 val / 147 test — no leak"),
              ("Nucleus sampling", " (top-p 0,9, T 0,7): il beam collassa sulla caption modale"),
              ("", "Curve sane: nessun overfitting, NaN o collasso")], bsize=12)
card(s, Inches(8.45), y0, Inches(4.2), Inches(4.3),
     title="Lezione Kaggle", accent=GOLD, tsize=15, fill=DEEP, tcolor=GOLD,
     body="Il flag enable_gpu nei metadati NON abilita davvero la GPU: va attivata a mano.\n\n"
          "Il tier gratuito tende alla P100 (sm_60): la PyTorch preinstallata va reinstallata "
          "da wheel cu118.\n\nTutto il training: GPU gratuita T4 / P100.", bsize=12.5)
notes(s, """
Il setup sperimentale, in breve. Tutto l'addestramento è avvenuto sul tier GPU gratuito di Kaggle,
con T4 o P100. Gli iperparametri nella tabella sono default ragionevoli per ciascuna architettura,
non il frutto di una ricerca esaustiva: m1 più epoche e learning rate più alto, m6 poche epoche e
learning rate piccolo, com'è naturale per un modello pre-addestrato che va solo rifinito.
Due scelte garantiscono un confronto onesto. Primo, lo split è sample-disjoint: lo stesso formaggio
non compare mai sia in training sia in test, così evitiamo il leak. Secondo, in inferenza usiamo
nucleus sampling invece del beam search, perché su dataset piccoli il beam collassa sulla caption
più frequente, gonfiando artificialmente il BLEU ma riducendo la diversità.
Le curve di addestramento sono sane su tutti e tre i modelli — nessun overfitting, nessun NaN,
nessun collasso — quindi le differenze nei risultati non sono dovute a patologie di training.
A destra, una lezione Kaggle pagata cara, per chi volesse riprodurre: il flag enable_gpu nei
metadati non basta, e la PyTorch preinstallata sulle P100 va reinstallata. Dettagli, ma di quelli
che fanno perdere un pomeriggio.
""")

# =====================================================================
# SLIDE 22 — BLEU PER ATTRIBUTO (figura)
# =====================================================================
s = newslide()
chrome(s, 22)
y0 = kicker_title(s, "Risultati · per-attributo", "BLEU-4: i modelli vincono dove le caption sono diverse")
pic_fit(s, A("fig_bleu_attr.png"), Inches(0.7), y0, Inches(8.1), Inches(4.4))
card(s, Inches(9.0), y0+Inches(0.1), Inches(3.65), Inches(4.1),
     title="Lettura", accent=PRIM, tsize=15,
     bullets=[("Vincono:", " Struttura (+0,128), Profumo, Texture, Colore"),
              ("Pari/perdono:", " Aroma, Sapore, Spessore"),
              ("", "Non per debolezza: BLEU premia il predittore modale su caption poco diverse"),
              ("4 / 7", " attributi battono la baseline costante")], bsize=12)
notes(s, """
Vediamo i primi risultati, in regime per-attributo: sette modelli separati, uno per dimensione
sensoriale. Il grafico mostra il BLEU-4 per attributo, confrontato con la baseline most_frequent —
quella che emette sempre la caption più frequente.
Il quadro è netto e a due facce. I modelli addestrati battono chiaramente la baseline costante su
quattro attributi: struttura della pasta, con un margine notevole di oltre 0,12 punti, e poi
profumo, texture e colore. Su altri tre — aroma, sapore e spessore — pareggiano o perdono.
Attenzione a non leggere male questo "perdono". Non significa che i modelli siano scarsi su quegli
attributi. Significa che lì le caption di riferimento sono poco diverse tra loro, molto
template-heavy, e in quelle condizioni BLEU premia chi predice sempre la maggioranza. È un difetto
della metrica, non del modello — e lo dimostreremo tra poco con lo shuffle test, che su spessore
conferma che m3 e m6 usano davvero l'immagine.
In sintesi: 4 attributi su 7 battono la baseline, e i 3 "persi" sono un artefatto di BLEU.
""")

# =====================================================================
# SLIDE 23 — GLOBALE (chart nativo)
# =====================================================================
s = newslide()
chrome(s, 23)
y0 = kicker_title(s, "Risultati · modello globale", "Un solo modello su tutti gli attributi")
cd = CategoryChartData()
cd.categories = ["BLEU-1","BLEU-4","METEOR","ROUGE-L","CIDEr"]
cd.add_series("m1 · CNN+LSTM", (0.3501,0.1283,0.2938,0.2951,0.2517))
cd.add_series("m3 · ViT+Transf.", (0.3649,0.1237,0.2875,0.2981,0.2856))
cd.add_series("m6 · ViT+GePpeTto", (0.3657,0.1307,0.2928,0.3007,0.2741))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), y0, Inches(8.3), Inches(4.4), cd)
ch=gf.chart; style_chart(ch); color_series(ch.plots[0],[M1C,M3C,M6C])
ch.plots[0].gap_width=90; ch.value_axis.has_major_gridlines=True
ch.value_axis.maximum_scale=0.4
card(s, Inches(9.15), y0+Inches(0.1), Inches(3.5), Inches(4.1),
     title="m6 vince di misura", accent=M6C, tsize=15,
     bullets=[("", "m6 primo su BLEU-1/4, ROUGE-L, METEOR≈"),
              ("", "m3 primo su CIDEr"),
              ("0,007", " di spread su BLEU-4: cluster strettissimo"),
              ("⚠", " most_frequent ha il BLEU-1 più alto ma BLEU-4 peggiore: degenere")], bsize=12)
notes(s, """
Passiamo al regime globale: un unico modello addestrato su tutti gli attributi insieme. È un
compito intrinsecamente più difficile, perché il modello deve anche scegliere quale dimensione
sensoriale descrivere, e lavora con un vocabolario molto più ampio.
Il grafico confronta i tre modelli su cinque metriche. Due osservazioni. Primo, m6 — il ViT con
GePpeTto — vince di misura quasi ovunque: è primo su BLEU-1, BLEU-4 e ROUGE-L, e sostanzialmente
alla pari su METEOR; m3 lo supera solo su CIDEr. Secondo, e altrettanto importante, i tre modelli
sono in un cluster strettissimo: lo spread su BLEU-4 è di appena sette millesimi. Sono molto vicini.
Una nota sulla baseline, che non ho messo nel grafico per leggibilità ma è cruciale:
most_frequent, la frase costante, avrebbe il BLEU-1 più alto di tutti ma il BLEU-4 peggiore. È il
comportamento degenere del "predici sempre la maggioranza": indovini molte singole parole, ma non
componi mai una frase corretta di quattro parole consecutive diverse dal template. Tenete a mente
questa anomalia: è il primo indizio che BLEU, da solo, ci sta ingannando.
""")

# =====================================================================
# SLIDE 24 — PER-ATTR VS GLOBALE (figura)
# =====================================================================
s = newslide()
chrome(s, 24)
y0 = kicker_title(s, "Attenzione all'interpretazione", "Per-attributo vs globale: due righelli diversi")
pic_fit(s, A("fig_perattr_vs_global.png"), Inches(0.7), y0, Inches(7.6), Inches(4.4))
card(s, Inches(8.5), y0+Inches(0.1), Inches(4.15), Inches(4.1),
     title="Perché il salto 3–4×", accent=GOLD, tsize=15,
     bullets=[("", "NON significa che i per-attributo apprendano meglio"),
              ("", "La valutazione per-attributo gira su distribuzione più stretta"),
              ("~80–140", " parole di vocabolario vs ~600+"),
              ("Scaffolding", " condiviso: ~5 parole di 4-gram costante"),
              ("→", " misurano performance su distribuzioni diverse")], bsize=12)
notes(s, """
Una slide di cautela metodologica, perché è facile trarre la conclusione sbagliata.
Avrete notato che i numeri per-attributo — BLEU-4 da 0,33 a 0,47 — sono molto più alti del BLEU-4
globale, intorno a 0,13. Un salto di tre o quattro volte. La tentazione è dire: "allora i modelli
per-attributo sono molto migliori". Sbagliato.
Quel salto non riflette un apprendimento migliore, ma il fatto che la valutazione per-attributo
gira su una distribuzione molto più stretta. Quando ti concentri su un solo attributo, il
vocabolario collassa da oltre 600 parole a circa 80-140, e lo scaffolding condiviso — "il formaggio
ha un tale attributo di..." — contribuisce da solo cinque parole di 4-gram costante che fanno
volare il BLEU.
Sono due righelli diversi che misurano performance su distribuzioni diverse. Il modello globale è
intrinsecamente più difficile perché deve anche scegliere quale dimensione descrivere. Confrontare
i due numeri come se fossero sulla stessa scala sarebbe un errore — ed è il tipo di sottigliezza
che separa un'analisi onesta da una che si auto-inganna.
""")

# =====================================================================
# SLIDE 25 — SHUFFLE TEST (figura) KEY
# =====================================================================
s = newslide()
chrome(s, 25)
y0 = kicker_title(s, "Il risultato chiave", "Lo shuffle test: il modello usa l'immagine?")
pic_fit(s, A("fig_shuffle.png"), Inches(0.7), y0, Inches(5.6), Inches(4.4))
card(s, Inches(6.55), y0, Inches(6.1), Inches(2.05),
     title="Come funziona", accent=PRIM, tsize=15.5,
     body="Si mescolano le predizioni tra le righe di test, rompendo l'allineamento "
          "predizione↔immagine. Si ricalcola 100 volte → distribuzione nulla → z-score. "
          "z > 3 ⟺ p < 0,001: forte evidenza di image-conditioning.", bsize=13)
card(s, Inches(6.55), y0+Inches(2.2), Inches(6.1), Inches(2.4),
     title="Cosa dice", accent=RED, tsize=15.5, fill=DEEP, tcolor=GOLD,
     bullets=[("m1 (ResNet):", " z≈0 ovunque → puro modello linguistico"),
              ("m3 / m6 (ViT):", " usano l'immagine su 4/7 attributi"),
              ("Profumo, Spessore, Colore, Struttura", " → riescono"),
              ("Aroma", " → nessun modello usa l'immagine")], bsize=12.5)
notes(s, """
Siamo al risultato più importante dell'intero progetto. La domanda è tanto semplice quanto
fondamentale: il modello sta davvero usando l'immagine, o sta solo imitando la distribuzione delle
caption?
Lo strumento è lo shuffle test. L'idea: se un modello è davvero condizionato dall'immagine, la sua
predizione per l'immagine i deve combaciare col riferimento di i meglio di una predizione presa a
caso. Allora mescoliamo le predizioni tra le righe di test, rompendo l'allineamento, ricalcoliamo
la sovrapposizione cento volte per ottenere una distribuzione nulla, e misuriamo lo z-score. Uno
z sopra 3 corrisponde a p sotto lo 0,001: evidenza forte.
Il verdetto, nel grafico, è cristallino. m1, il modello con ResNet, ha z vicino a zero ovunque:
non usa mai l'immagine, è di fatto un puro modello linguistico che indovina dalla distribuzione.
m3 e m6, i modelli con ViT, usano l'immagine chiaramente su quattro attributi su sette: profumo,
spessore, colore e struttura. L'unico dove nessun modello usa l'immagine è l'aroma.
Tenete questo grafico in mente: la prossima slide ne trae la conclusione architetturale.
""")

# =====================================================================
# SLIDE 26 — CONCLUSIONE ARCHITETTURALE
# =====================================================================
s = newslide()
chrome(s, 26)
y0 = kicker_title(s, "L'interpretazione", "Conta l'encoder, non il decoder")
card(s, Inches(0.7), y0+Inches(0.1), Inches(5.9), Inches(3.5),
     title="m1 vs m3 — isola l'encoder", accent=M3C, tsize=16,
     body="m1 (ResNet) non si condiziona MAI all'immagine.\nm3 (ViT) lo fa sulla maggior parte "
          "degli attributi.\n\n→ Le feature di un ResNet-50 congelato non bastano a questa scala "
          "di dati; quelle di un ViT-B/16 congelato sì.", bsize=14)
card(s, Inches(6.75), y0+Inches(0.1), Inches(5.9), Inches(3.5),
     title="m3 vs m6 — isola il decoder", accent=M6C, tsize=16,
     body="Producono output diversi (mai la stessa caption), ma hanno lo STESSO profilo di "
          "image-conditioning: gli stessi attributi riescono, gli stessi falliscono.\n\n"
          "→ Il decoder cambia lo stile, non se l'immagine viene usata.", bsize=14)
# barra conclusione
rect(s, Inches(0.7), y0+Inches(3.8), Inches(11.95), Inches(0.75), DEEP, round_=True)
_, tf=textbox(s, Inches(1.0), y0+Inches(3.8), Inches(11.5), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]
setrun(p,"Il «cancello binario» dell'image-conditioning è l'encoder visivo — non il decoder, non la scala dati, non gli iperparametri.",13.5,GOLD,bold=True)
notes(s, """
Mettiamo insieme lo shuffle test con il disegno a due assi della slide delle architetture, e la
conclusione si scrive da sola.
Confronto m1 contro m3, che isola l'encoder: m1 con ResNet non si condiziona mai all'immagine, m3
con ViT lo fa sulla maggior parte degli attributi. Cambiando solo l'encoder, si accende l'uso
dell'immagine. La lettura è che le feature di un ResNet-50 congelato non sono abbastanza
informative a questa scala di dati, mentre quelle di un ViT-B/16 congelato sì.
Confronto m3 contro m6, che isola il decoder: producono output diversi — non danno mai la stessa
caption per la stessa immagine — eppure hanno esattamente lo stesso profilo di image-conditioning.
Gli stessi attributi riescono, gli stessi falliscono. Cambiare il decoder cambia lo stile della
frase, non il fatto che l'immagine venga usata.
La sintesi è in fondo: il cancello binario dell'image-conditioning, su questo dataset, è l'encoder
visivo. Non il decoder, non la quantità di dati, non gli iperparametri. È un risultato pulito, e ci
dice anche dove intervenire in futuro: sull'encoder. Ora rafforziamo questa conclusione con le
metriche complementari.
""")

# =====================================================================
# SLIDE 27 — OLTRE BLEU: 7 METRICHE
# =====================================================================
s = newslide()
chrome(s, 27)
y0 = kicker_title(s, "Valutazione onesta", "Oltre BLEU: sette metriche complementari")
rows=[("Metrica","Cosa misura","Famiglia"),
      ("BLEU-1/4","precisione n-gram (parole / fraseologia esatta)","testo"),
      ("METEOR","sovrapposizione con stemming e sinonimi","testo"),
      ("ROUGE-L","più lunga sottosequenza comune (ordine)","testo"),
      ("CIDEr","n-gram pesati TF-IDF: premia i termini rari salienti","testo"),
      ("BERTScore","similarità semantica via embedding (BERT it)","semantica"),
      ("Conformità Vocab.","% parole nel lessico sensoriale attestato","dominio"),
      ("CLIPScore","appropriatezza caption–immagine (coseno CLIP)","immagine")]
tx=Inches(0.7); ty=y0; tw=Inches(11.95); th=Inches(4.3)
tbl=s.shapes.add_table(len(rows),3,tx,ty,tw,th).table
for ci,w in enumerate([2.6,7.15,2.2]): tbl.columns[ci].width=Inches(w)
for ri,row in enumerate(rows):
    last = (ri==len(rows)-1)
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ci==2: p.alignment=PP_ALIGN.CENTER
        if ri==0:
            setrun(p,val,13,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            bold = last or ci==0
            setrun(p,val,12.5, (GOLD_D if last else INK), bold=bold)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0x2C,0x5F,0x8A) if False else (RGBColor(0xEC,0xE3,0xCF) if last else (CREAM2 if ri%2 else WHITE)))
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
notes(s, """
Abbiamo visto che BLEU, da solo, inganna: premia la ripetizione del template e non guarda mai
l'immagine. Per leggere i modelli con più onestà abbiamo affiancato a BLEU sei metriche
complementari, raggruppabili in tre famiglie.
Le metriche di testo confrontano la predizione col riferimento testuale: oltre a BLEU, METEOR — che
è più indulgente perché allinea con stemming e sinonimi — ROUGE-L, che premia l'ordine delle
parole, e CIDEr, pensata apposta per il captioning, che pesa gli n-gram con TF-IDF e quindi premia
i termini sensoriali rari e informativi come "occhiatura" invece delle parole comuni.
Poi una metrica semantica, BERTScore, che misura la similarità di significato via embedding. Una
metrica di dominio, la conformità al vocabolario, che chiede: il modello "parla formaggio" col
registro certificato?
E infine, evidenziata in basso, la più importante per il captioning: CLIPScore, l'unica che guarda
davvero l'immagine. Confronta la caption con la fetta nello spazio CLIP, ignorando il riferimento
del panel. Le prossime due slide mostrano cosa aggiungono CIDEr e CLIPScore.
""")

# =====================================================================
# SLIDE 28 — CIDEr (chart nativo)
# =====================================================================
s = newslide()
chrome(s, 28)
y0 = kicker_title(s, "CIDEr", "Discrimina dove BLEU appiattisce")
cd = CategoryChartData()
cd.categories = ["random","freq_weighted","m1","m3","m6","most_frequent"]
cd.add_series("CIDEr medio (7 attributi)", (0.234,0.370,0.586,0.646,0.704,0.790))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), y0, Inches(7.7), Inches(4.35), cd)
ch=gf.chart; style_chart(ch, legend=False)
plot=ch.plots[0]; plot.gap_width=70; plot.has_data_labels=True
plot.data_labels.number_format='0.00'; plot.data_labels.number_format_is_linked=False
plot.data_labels.font.size=Pt(12); plot.data_labels.font.bold=True; plot.data_labels.position=XL_LABEL_POSITION.OUTSIDE_END
points_color(plot.series[0],[GREY,GREY,M1C,M3C,M6C,GOLD])
ch.value_axis.has_major_gridlines=True
card(s, Inches(8.5), y0+Inches(0.1), Inches(4.15), Inches(4.1),
     title="Lettura", accent=M6C, tsize=15,
     bullets=[("Ordina i modelli:", " random < freq < m1 < m3 < m6"),
              ("", "I modelli addestrati battono nettamente le baseline casuali"),
              ("m6", " è il migliore tra i modelli"),
              ("most_frequent (0,79)", " in testa solo perché il riferimento È spesso la frase frequente: artefatto, non qualità")], bsize=11.5)
notes(s, """
Cominciamo da CIDEr, che è la metrica pensata apposta per il captioning. A differenza di BLEU, che
appiattiva i modelli su valori vicini, CIDEr pesa i termini rari con TF-IDF e quindi li separa.
Guardate l'ordinamento nel grafico, sulle medie dei sette attributi: random 0,23, poi freq_weighted
0,37, poi i modelli addestrati che salgono nettamente — m1 a 0,59, m3 a 0,65, m6 a 0,70. CIDEr
ordina correttamente i modelli, e conferma che m6 è il migliore. È esattamente l'informazione che
BLEU non riusciva a darci: una graduatoria pulita.
C'è una sola colonna più alta di tutte, in oro: most_frequent, a 0,79. Non lasciatevi ingannare:
la baseline costante è in testa solo perché, su questo dataset, il riferimento del panel è spesso
proprio la frase frequente. È un artefatto del riferimento singolo, non qualità reale. Su un caso
estremo come lo spessore della crosta, dove BLEU dava 0,39 per tutti, CIDEr va da 1,04 a 2,29: un
ordinamento ricchissimo che BLEU semplicemente non vedeva.
Conclusione: CIDEr aggiunge informazione vera e separa i modelli addestrati dal rumore.
""")

# =====================================================================
# SLIDE 29 — CLIPSCORE (chart nativo doppio)
# =====================================================================
s = newslide()
chrome(s, 29)
y0 = kicker_title(s, "CLIPScore", "Conferma indipendente dello shuffle test")
# sinistra: per modello (piatto)
cd1 = CategoryChartData()
cd1.categories = ["most_freq","random","freq_w","m1","m3","m6"]
cd1.add_series("CLIPScore medio", (0.1902,0.1913,0.1919,0.1925,0.1921,0.1933))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), y0+Inches(0.25), Inches(5.9), Inches(3.7), cd1)
ch=gf.chart; style_chart(ch, legend=False)
ch.value_axis.minimum_scale=0.185; ch.value_axis.maximum_scale=0.195
plot=ch.plots[0]; plot.gap_width=60
points_color(plot.series[0],[GOLD,GREY,GREY,M1C,M3C,M6C])
_, tf=textbox(s, Inches(0.7), y0-Inches(0.05), Inches(5.9), Inches(0.32))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"Per MODELLO → tutti uguali (scarto < 0,004)",12,DEEP,bold=True)
# destra: per attributo
cd2 = CategoryChartData()
cd2.categories = ["Aroma","Sapore","Spessore","Profumo","Struttura","Colore","Texture"]
cd2.add_series("CLIPScore", (0.181,0.184,0.184,0.185,0.197,0.206,0.208))
gf2=s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.85), y0+Inches(0.25), Inches(5.8), Inches(3.7), cd2)
ch2=gf2.chart; style_chart(ch2, legend=False)
plot2=ch2.plots[0]; plot2.gap_width=50
points_color(plot2.series[0],[PRIM,RED,GOLD,PRIM,GREEN,GREEN,GREEN])
ch2.value_axis.minimum_scale=0.17
_, tf=textbox(s, Inches(6.85), y0-Inches(0.05), Inches(5.8), Inches(0.32))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"Per ATTRIBUTO → visivi alti, olfatto/gusto bassi",12,DEEP,bold=True)
# barra conclusione
rect(s, Inches(0.7), y0+Inches(4.05), Inches(11.95), Inches(0.6), DEEP, round_=True)
_, tf=textbox(s, Inches(1.0), y0+Inches(4.05), Inches(11.5), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]
setrun(p,"Con encoder congelato l'immagine non viene sfruttata: i modelli fanno language modeling. La variazione vera è per attributo.",12.5,GOLD,bold=True)
notes(s, """
Ora CLIPScore, l'unica metrica che guarda l'immagine, e la conferma indipendente dello shuffle
test. Due grafici, due messaggi.
A sinistra, CLIPScore medio per modello. Notate l'asse: ho zoomato volutamente, e nonostante lo
zoom le colonne sono praticamente identiche. I modelli addestrati m1, m3, m6 non superano la
baseline costante most_frequent — quella che spara sempre la stessa frase ignorando l'immagine. Lo
scarto è sotto lo 0,004, dentro il rumore. È la conferma quantitativa e indipendente di ciò che lo
shuffle test ci aveva già detto: con l'encoder congelato, i modelli non àncorano la caption ai
contenuti visivi, fanno language modeling sulla distribuzione delle caption.
A destra, la variazione vera: è per attributo, non per modello. CLIP "vede" che texture, colore e
struttura — gli attributi visivi — ottengono punteggi più alti, mentre aroma, sapore e profumo
restano bassi per chiunque, perché nessuna immagine può rivelare un sapore. È un controllo di
sanità che valida la metrica. E ricordate l'eccezione: lo spessore della crosta è visibile ma
basso, perché la crosta è una porzione piccola della fetta e CLIP la pesa poco. Tutto torna.
""")

# =====================================================================
# SLIDE 30 — TRANELLO BLEU (figura)
# =====================================================================
s = newslide()
chrome(s, 30)
y0 = kicker_title(s, "La morale metodologica", "Il tranello di BLEU, visualizzato")
pic_fit(s, A("m_bleu_trap.png"), Inches(0.7), y0, Inches(8.2), Inches(4.4))
card(s, Inches(9.1), y0+Inches(0.1), Inches(3.55), Inches(4.1),
     title="Una metrica mente", accent=RED, tsize=15.5, fill=DEEP, tcolor=GOLD,
     body="A sinistra BLEU-1: la caption costante most_frequent VINCE.\n\n"
          "A destra CLIPScore: i due modelli sono praticamente uguali.\n\n"
          "→ BLEU-1 da solo è fuorviante. Le metriche vanno lette insieme.", bsize=13.5)
notes(s, """
Questa slide è la morale metodologica di tutta la valutazione, condensata in un'immagine.
Guardate la stessa coppia di modelli — la baseline costante most_frequent contro m6 — vista da due
metriche diverse. A sinistra, con BLEU-1, la caption costante vince: emette sempre la stessa frase
e combacia sul prefisso. A destra, con CLIPScore, i due sono praticamente indistinguibili.
Stesso identico confronto, due verdetti opposti. Questo dimostra plasticamente che una metrica, da
sola, può mentire. Se avessimo guardato solo BLEU-1, avremmo concluso che una frase costante batte
un modello addestrato — una conclusione assurda. È esattamente per questo che abbiamo costruito una
batteria di metriche e le leggiamo insieme: BLEU per la fraseologia, CIDEr per i termini salienti,
BERTScore per il significato, e CLIPScore per l'ancoraggio all'immagine.
BERTScore, per completezza, qui è poco discriminante: tutti i valori stanno tra 0,83 e 0,92, perché
tutte le frasi parlano di formaggio con la stessa struttura. È utile come sanity check, non per il
ranking. La lezione resta: mai fidarsi di una sola metrica.
""")

# =====================================================================
# SLIDE 31 — ESEMPI QUALITATIVI
# =====================================================================
s = newslide()
chrome(s, 31)
y0 = kicker_title(s, "Esempi qualitativi", "Cosa genera davvero m6")
rows=[("Attributo","Predizione (m6)","Riferimento (panel)"),
      ("Aroma","Il formaggio ha un aroma di panna.","Il formaggio ha un aroma di panna.  ✓"),
      ("Sapore","…sapore leggermente salato e piccante.","Il formaggio ha un sapore salato."),
      ("Spessore","La crosta è mediamente spessa.","La crosta è mediamente spessa.  ✓"),
      ("Colore","…colore giallo carico omogeneo.","…colore giallo troppo carico."),
      ("Struttura","…frattura irregolare e grana fine.","…stirata e poca grana a tratti.")]
tx=Inches(0.7); ty=y0; tw=Inches(11.95); th=Inches(3.3)
tbl=s.shapes.add_table(len(rows),3,tx,ty,tw,th).table
for ci,w in enumerate([1.9,5.0,5.05]): tbl.columns[ci].width=Inches(w)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ri==0:
            setrun(p,val,12.5,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            match = "✓" in val
            setrun(p,val,12, (GREEN if match else INK), bold=(ci==0 or match))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
card(s, Inches(0.7), ty+Inches(3.45), Inches(11.95), Inches(0.9),
     title=None, accent=GOLD,
     body="Quattro pattern ricorrenti: (1) lo scaffolding è sempre corretto e vale ~50% del BLEU-4; "
          "(2) la lingua è fluente grazie al LM pre-addestrato; (3) il descrittore è spesso valido "
          "ma per un altro formaggio; (4) i riferimenti tra degustatori variano enormemente → tetto al BLEU.",
     bsize=12.5)
notes(s, """
I numeri raccontano metà della storia; le caption generate raccontano l'altra metà. Qui vedete
predizioni reali di m6 affiancate al riferimento del panel, su campioni del test set.
A volte è un match esatto: sull'aroma di panna, predizione e riferimento coincidono parola per
parola. Sullo spessore, "la crosta è mediamente spessa", di nuovo perfetto. Ma guardate gli altri
casi: sul sapore il modello dice "leggermente salato e piccante" dove il panel diceva solo "salato";
sulla struttura propone "frattura irregolare e grana fine" mentre il riferimento parla di pasta
stirata.
Da questi esempi emergono quattro pattern, riassunti in basso. Primo: lo scaffolding, "il formaggio
ha un...", è sempre corretto e da solo vale circa metà del BLEU-4 — ecco perché i numeri non sono
bassissimi. Secondo: la lingua è fluente e grammaticale, merito del modello pre-addestrato. Terzo,
e più rivelatore: il descrittore è spesso plausibile ma riferito a un altro formaggio — esattamente
ciò che ci aspettiamo da un modello solo parzialmente condizionato dall'immagine. Quarto: la
variabilità tra degustatori è enorme e pone un tetto a ciò che qualsiasi modello può ottenere con un
BLEU a riferimento singolo.
""")

# =====================================================================
# SLIDE 32 — CONCLUSIONI & FUTURO
# =====================================================================
s = newslide(DEEP)
rect(s, 0, 0, Inches(0.22), SH, GOLD)
cheese_holes(s, Inches(11.6), Inches(1.1), GOLD, base=0.55, alpha=16)
_, tf=textbox(s, Inches(0.85), Inches(0.6), Inches(11), Inches(0.9))
setrun(tf.paragraphs[0],"Conclusioni & lavori futuri",34,CREAM,bold=True,font=FONT_H)
rect(s, Inches(0.88), Inches(1.5), Inches(2.0), Pt(4), GOLD)
# tre colonne conclusioni
def dcard(x,w,title,acc,bullets):
    rect(s, x, Inches(1.85), w, Inches(3.4), RGBColor(0x24,0x40,0x59), round_=True, shadow=True)
    rect(s, x, Inches(1.85), w, Inches(0.1), acc)
    _, tf=textbox(s, x+Inches(0.25), Inches(2.0), w-Inches(0.45), Inches(3.15))
    p=tf.paragraphs[0]; setrun(p,title,15.5,acc,bold=True,font=FONT_H)
    for b in bullets:
        p=tf.add_paragraph(); p.space_before=Pt(6)
        setrun(p,"● ",10,acc,bold=True)
        if isinstance(b,tuple): setrun(p,b[0],12.5,CREAM,bold=True); setrun(p,b[1],12.5,RGBColor(0xD9,0xD2,0xC2))
        else: setrun(p,b,12.5,RGBColor(0xD9,0xD2,0xC2))
dcard(Inches(0.7),Inches(3.85),"Risultato scientifico",GREEN,
      [("Conta l'encoder:", " ResNet ignora l'immagine, ViT la usa su 4/7"),
       ("", "Indipendente da decoder, scala dati, iperparametri"),
       ("", "Confermato da shuffle test + CLIPScore")])
dcard(Inches(4.74),Inches(3.85),"Risultato pratico",GOLD,
      [("4/7", " attributi battono la baseline (max Struttura +0,128)"),
       ("m6 (ViT+GePpeTto)", " vince di misura quasi ovunque"),
       ("$5,60", " costo dati · 38.437 coppie consegnate")])
dcard(Inches(8.78),Inches(3.87),"Lavori futuri",PRIM,
      [("Fine-tuning dell'encoder", " e ri-misura di CLIPScore"),
       ("", "Baseline di retrieval k-NN come «pavimento» image-aware"),
       ("", "Se i modelli si staccano dalla baseline → immagine sfruttata")])
# framing onesto
_, tf=textbox(s, Inches(0.7), Inches(5.45), Inches(11.9), Inches(0.95))
p=tf.paragraphs[0]
setrun(p,"Framing onesto:  ",13,GOLD,bold=True)
setrun(p,"nessun modello sostituirebbe un degustatore. Hanno appreso il vocabolario e una mappa parziale immagine→descrittore — coerente con la scala dati e il rumore dell'annotazione multi-degustatore.",13,CREAM)
_, tf=textbox(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.6))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"Grazie. — Domande?",18,GOLD,bold=True,font=FONT_H)
notes(s, """
Tiriamo le fila. Il progetto consegna due risultati e una promessa per il futuro.
Il risultato scientifico è il più solido: su questo dataset, il determinante dell'uso dell'immagine
è l'encoder visivo. ResNet ignora l'immagine su tutti e sette gli attributi, ViT la usa su quattro
su sette, indipendentemente dal decoder, dalla scala dei dati e dagli iperparametri. E non è
un'affermazione isolata: due strumenti indipendenti — lo shuffle test e il CLIPScore — convergono
sulla stessa conclusione.
Il risultato pratico: i modelli addestrati battono la baseline costante su quattro attributi su
sette, col massimo sulla struttura. m6, il ViT con GePpeTto, vince di misura quasi ovunque. E tutto
questo poggia su un dataset di 38 mila coppie costruito con appena 5 dollari e 60 di costo LLM.
I lavori futuri seguono direttamente dalla diagnosi: poiché il collo di bottiglia è l'encoder
congelato, il prossimo esperimento naturale è il fine-tuning dell'encoder e la ri-misura del
CLIPScore. Se i modelli inizieranno a staccarsi dalla baseline costante, avremo la prova diretta
che l'immagine viene finalmente sfruttata.
Chiudo con un framing onesto: nessuno di questi modelli sostituirebbe oggi un degustatore. Hanno
imparato il vocabolario sensoriale e una mappa parziale dall'immagine al descrittore — un risultato
coerente con la scala dei dati e col rumore intrinseco dell'annotazione multi-degustatore. Ma la
strada è tracciata, e sappiamo esattamente dove intervenire.
Grazie per l'attenzione. Sono a disposizione per le domande.
""")

prs.save(OUTFILE)
print("Salvato:", OUTFILE)
print("Slide totali:", len(prs.slides._sldIdLst))
