# -*- coding: utf-8 -*-
"""
Builds the presentation (English): Sensory captioning of Trentingrana cheese.
~32 slides, modern style with cheese cues, native charts + figures + real photos,
speaker notes for a ~30 minute talk.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from PIL import Image

ROOT = r"C:\Users\nicol\Desktop\Cheese"
ASSETS = os.path.join(ROOT, "presentation", "assets")
OUTFILE = os.path.join(ROOT, "presentation", "Trentingrana_Captioning_EN.pptx")

# ----------------- PALETTE (cheese cues) -----------------
CREAM   = RGBColor(0xF7, 0xF1, 0xE3)
CREAM2  = RGBColor(0xFB, 0xF7, 0xEE)
INK     = RGBColor(0x2B, 0x2A, 0x26)
PRIM    = RGBColor(0x2C, 0x5F, 0x8A)
DEEP    = RGBColor(0x1B, 0x33, 0x4A)
GOLD    = RGBColor(0xE0, 0xA4, 0x58)
GOLD_D  = RGBColor(0xC8, 0x88, 0x36)
GREEN   = RGBColor(0x5B, 0x8C, 0x5A)
RED     = RGBColor(0xC1, 0x49, 0x2E)
GREY    = RGBColor(0x8A, 0x8A, 0x84)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
M1C     = RGBColor(0xE0, 0x7A, 0x5F)
M3C     = RGBColor(0x3D, 0x70, 0x68)
M6C     = RGBColor(0x8A, 0x50, 0x82)

FONT_H  = "Georgia"
FONT_B  = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ----------------- HELPERS -----------------
def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _no_autosize(tf):
    tf.word_wrap = True

def add_bg(slide, color=CREAM):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, x, y, w, h, color, line=None, lw=None, shadow=False, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_fill(shp, color)
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(lw or 1)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        sp = el.makeelement(qn('a:effectLst'), {})
        outer = el.makeelement(qn('a:outerShdw'),
            {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val':'1B334A'})
        alpha = el.makeelement(qn('a:alpha'), {'val':'24000'})
        clr.append(alpha); outer.append(clr); sp.append(outer); el.append(sp)
    return shp

def oval(slide, x, y, w, h, color, alpha=None, line=None, lw=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    _set_fill(shp, color); shp.line.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(lw or 1)
    shp.shadow.inherit = False
    if alpha is not None:
        sp = shp.fill.fore_color._xFill
        srgb = sp.find(qn('a:srgbClr'))
        a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha*1000))})
        srgb.append(a)
    return shp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; _no_autosize(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tb, tf

def setrun(p, text, size, color=INK, bold=False, italic=False, font=FONT_B, spacing=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    if spacing is not None:
        rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(spacing))
    return r

def para(tf, first=False):
    if first and not tf.paragraphs[0].runs and tf.paragraphs[0].text == "":
        return tf.paragraphs[0]
    return tf.add_paragraph()

def cheese_holes(slide, cx, cy, color=GOLD, base=0.42, n=5, alpha=None):
    spots = [(0,0,1.0),(0.9,0.5,0.55),(-0.7,0.7,0.45),(0.5,-0.8,0.5),(-0.9,-0.4,0.4),(1.1,-0.5,0.35)]
    for i,(dx,dy,sc) in enumerate(spots[:n]):
        d = Inches(base*sc)
        oval(slide, cx+Inches(dx*0.8)-d/2, cy+Inches(dy*0.8)-d/2, d, d, color, alpha=alpha)

PROJ = "AI4FQC · Project 07 — GRANA Captioning"
def chrome(slide, idx, accent=GOLD, dark=False):
    rect(slide, 0, 0, Inches(0.16), SH, accent)
    fc = CREAM if dark else RGBColor(0x6B,0x68,0x60)
    _, tf = textbox(slide, Inches(0.45), SH-Inches(0.42), Inches(9), Inches(0.32))
    setrun(tf.paragraphs[0], PROJ, 9.5, fc, italic=True)
    _, tf2 = textbox(slide, SW-Inches(1.4), SH-Inches(0.42), Inches(0.95), Inches(0.32))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    setrun(p2, f"{idx:02d}", 11, accent, bold=True)

def kicker_title(slide, kicker, title, tcolor=DEEP, kcolor=GOLD_D, top=Inches(0.55)):
    _, tf = textbox(slide, Inches(0.7), top, Inches(11.9), Inches(0.4))
    setrun(tf.paragraphs[0], kicker.upper(), 12.5, kcolor, bold=True, spacing=220)
    _, tf2 = textbox(slide, Inches(0.7), top+Inches(0.42), Inches(11.9), Inches(1.0))
    setrun(tf2.paragraphs[0], title, 31, tcolor, bold=True, font=FONT_H)
    rect(slide, Inches(0.72), top+Inches(1.32), Inches(1.5), Pt(3), GOLD)
    return top+Inches(1.55)

def pic_fit(slide, path, box_x, box_y, box_w, box_h, align="center", valign="middle"):
    with Image.open(path) as im:
        iw, ih = im.size
    bw, bh = box_w, box_h
    ar = iw/ih; bar = bw/bh
    if ar > bar: w = bw; h = int(bw/ar)
    else: h = bh; w = int(bh*ar)
    if align == "center": x = box_x + (bw-w)//2
    elif align == "left": x = box_x
    else: x = box_x + (bw-w)
    if valign == "middle": y = box_y + (bh-h)//2
    elif valign == "top": y = box_y
    else: y = box_y + (bh-h)
    return slide.shapes.add_picture(path, x, y, width=w, height=h)

def A(name): return os.path.join(ASSETS, name)

def style_chart(chart, font_sz=11, legend=True, legpos=XL_LEGEND_POSITION.BOTTOM):
    chart.has_title = False
    if legend:
        chart.has_legend = True
        chart.legend.position = legpos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(font_sz); chart.legend.font.name = FONT_B
    else:
        chart.has_legend = False
    for ax in (chart.category_axis, chart.value_axis):
        try:
            ax.tick_labels.font.size = Pt(font_sz)
            ax.tick_labels.font.name = FONT_B
            ax.tick_labels.font.color.rgb = INK
        except Exception:
            pass

def color_series(plot, colors):
    for s, c in zip(plot.series, colors):
        s.format.fill.solid(); s.format.fill.fore_color.rgb = c
        s.format.line.fill.background()

def points_color(series, colors):
    for pt, c in zip(series.points, colors):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
        pt.format.line.fill.background()

def card(slide, x, y, w, h, title=None, body=None, accent=PRIM, fill=CREAM2,
         tsize=15, bsize=12.5, bullets=None, tcolor=None):
    rect(slide, x, y, w, h, fill, round_=True, shadow=True)
    rect(slide, x, y, Inches(0.09), h, accent, round_=False)
    dark = (fill == DEEP)
    txt = CREAM if dark else INK
    txt2 = RGBColor(0xD9,0xD2,0xC2) if dark else INK
    pad = Inches(0.28)
    _, tf = textbox(slide, x+pad, y+Inches(0.16), w-pad-Inches(0.18), h-Inches(0.3))
    if title:
        for j, ln in enumerate(title.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            setrun(p, ln, tsize, tcolor or (GOLD if dark else DEEP), bold=True, font=FONT_H)
    if body:
        p = para(tf, first=(title is None))
        setrun(p, body, bsize, txt)
    if bullets:
        for i,(b) in enumerate(bullets):
            p = para(tf, first=(title is None and i==0 and not body))
            p.space_before = Pt(4)
            setrun(p, "● ", 10, accent, bold=True)
            if isinstance(b, tuple):
                setrun(p, b[0], bsize, txt, bold=True)
                setrun(p, b[1], bsize, txt2)
            else:
                setrun(p, b, bsize, txt)
    return

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()

def newslide(bg=CREAM):
    s = prs.slides.add_slide(BLANK); add_bg(s, bg); return s

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = newslide(DEEP)
pic_fit(s, A("cheese_fetta_1_sq.jpg"), Inches(8.4), Inches(0), Inches(4.93), SH)
veil = rect(s, Inches(8.0), 0, Inches(1.4), SH, DEEP)
sp = veil.fill.fore_color._xFill.find(qn('a:srgbClr'))
sp.append(sp.makeelement(qn('a:alpha'), {'val':'62000'}))
cheese_holes(s, Inches(1.1), Inches(1.0), GOLD, base=0.5, alpha=22)
rect(s, 0, 0, Inches(0.22), SH, GOLD)
_, tf = textbox(s, Inches(0.85), Inches(1.7), Inches(7.4), Inches(0.5))
setrun(tf.paragraphs[0], "AI4FQC · PROJECT 07", 14, GOLD, bold=True, spacing=300)
_, tf = textbox(s, Inches(0.8), Inches(2.25), Inches(7.6), Inches(2.6))
p=tf.paragraphs[0]; setrun(p, "Sensory captioning", 44, CREAM, bold=True, font=FONT_H)
p=tf.add_paragraph(); setrun(p, "of Trentingrana cheese", 44, GOLD, bold=True, font=FONT_H)
_, tf = textbox(s, Inches(0.85), Inches(4.7), Inches(7.0), Inches(1.0))
setrun(tf.paragraphs[0],
       "From dataset construction to comparing three encoder–decoder methods",
       17, RGBColor(0xD9,0xD2,0xC2), italic=True)
rect(s, Inches(0.88), Inches(5.55), Inches(2.0), Pt(3), GOLD)
_, tf = textbox(s, Inches(0.85), Inches(5.8), Inches(7), Inches(0.8))
p=tf.paragraphs[0]
setrun(p, "Technical report  ·  ", 13, RGBColor(0xC9,0xC2,0xB2))
setrun(p, "38,437 image–caption pairs  ·  3 models  ·  7 attributes", 13, CREAM, bold=True)
notes(s, """
Good morning everyone, and welcome. Today I'm presenting our project on the sensory captioning of
Trentingrana cheese, developed within the AI4FQC initiative — Artificial Intelligence for Food
Quality Control, Project 07.
The starting idea is simple to state but surprisingly rich to deliver: automatically generate
sensory descriptions, in Italian, from images of grana cheese cross-sections. In other words,
teach a model to "tell the story" of a cheese by looking at a slice of it.
The work has two pillars: building the dataset — the most original and demanding part, to which
I'll devote roughly 70% of the time — and comparing three conceptually different encoder-decoder
architectures. We'll look at concrete numbers: 38 thousand image-caption pairs, three models,
seven sensory attributes, and above all a clear scientific result about what truly drives a model
to actually use the image.
Let's start with the context.
""")

# =====================================================================
# SLIDE 2 — CONTEXT
# =====================================================================
s = newslide()
chrome(s, 2)
y0 = kicker_title(s, "The context", "Why an AI that describes cheese")
card(s, Inches(0.7), y0, Inches(6.5), Inches(1.7),
     title="Sensory quality control",
     body="Every wheel of grana is assessed by a panel of expert tasters. It is precious work, "
          "but it is slow, subjective and hard to scale across thousands of wheels.",
     accent=PRIM)
card(s, Inches(0.7), y0+Inches(1.85), Inches(6.5), Inches(1.95),
     title="The idea", accent=GOLD,
     bullets=[("Input: ", "IRIS images of cross-sections (controlled lighting)"),
              ("Output: ", "a sensory description in Italian, per attribute"),
              ("Value: ", "objective, reproducible support to the human panel")])
fx = Inches(7.55); fw = Inches(5.1)
rect(s, fx-Inches(0.1), y0-Inches(0.05), fw+Inches(0.2), Inches(3.9), WHITE, round_=True, shadow=True)
pic_fit(s, A("cheese_grana_1_sq.jpg"), fx, y0+Inches(0.08), fw, Inches(3.4))
_, tf = textbox(s, fx, y0+Inches(3.5), fw, Inches(0.35))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p, "Close-up of the grain — IRIS electronic visual analyzer", 10.5, GREY, italic=True)
notes(s, """
Let's begin with the real problem. Grana quality is certified by tasting panels: experts who
smell, taste and observe each wheel, describing aroma, flavor, structure and so on. It is a wealth
of expertise, but it has three limits: it is slow, it is subjective — two panelists describe the
same wheel with different words — and it is hard to scale when there are thousands of wheels.
The project's idea is to support this work with an automatic system. We feed the model an image
captured with the IRIS electronic visual analyzer, under controlled lighting, like the one you see
on the right, and we ask it to produce a sensory description in Italian for each attribute.
Important: the goal is not to replace the taster, but to offer objective, reproducible support. As
we'll see, some attributes — the visible ones — lend themselves far better than others, and this
very distinction will be one of the most interesting results.
""")

# =====================================================================
# SLIDE 3 — 7 ATTRIBUTES
# =====================================================================
s = newslide()
chrome(s, 3)
y0 = kicker_title(s, "The subject", "Seven sensory attributes")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.6))
setrun(tf.paragraphs[0],
       "For each wheel the panel describes seven dimensions. Three are visible in the image, four are not — a distinction that returns in the results.",
       13.5, INK)
attrs = [("Paste Colour","visual",GREEN),("Texture","visual",GREEN),
         ("Paste Structure","visual",GREEN),("Rind Thickness","visual*",GOLD_D),
         ("Smell","olfactory",PRIM),("Aroma","olfactory",PRIM),("Taste","gustatory",RED)]
cols=4; cw=Inches(2.92); ch=Inches(1.55); gx=Inches(0.18); gy=Inches(0.25)
x0=Inches(0.7); yA=y0+Inches(0.75)
for i,(name,kind,col) in enumerate(attrs):
    r=i//cols; c=i%cols
    x=x0+c*(cw+gx); y=yA+r*(ch+gy)
    rect(s, x, y, cw, ch, CREAM2, round_=True, shadow=True)
    rect(s, x, y, cw, Inches(0.12), col, round_=False)
    oval(s, x+cw-Inches(0.5), y+Inches(0.32), Inches(0.28), Inches(0.28), col, alpha=28)
    _, tf = textbox(s, x+Inches(0.22), y+Inches(0.28), cw-Inches(0.4), ch-Inches(0.4))
    p=tf.paragraphs[0]; setrun(p, name, 15, DEEP, bold=True, font=FONT_H)
    p=tf.add_paragraph(); p.space_before=Pt(4); setrun(p, kind.upper(), 11, col, bold=True, spacing=150)
x=x0+3*(cw+gx); y=yA+1*(ch+gy)
rect(s, x, y, cw, ch, DEEP, round_=True, shadow=True)
_, tf = textbox(s, x+Inches(0.22), y+Inches(0.2), cw-Inches(0.4), ch-Inches(0.4))
p=tf.paragraphs[0]; setrun(p, "* visible but small", 12, GOLD, bold=True)
p=tf.add_paragraph(); p.space_before=Pt(3)
setrun(p, "The rind is a small area: CLIP weights it little.", 11.5, CREAM)
notes(s, """
These are the seven sensory dimensions the panel rates for each wheel: paste colour, texture,
paste structure, rind thickness, smell, aroma and taste.
Please keep in mind a distinction that will become central in the second part. Three attributes —
paste colour, texture and paste structure — are intrinsically visual: they can be read from the
image. Three — smell, aroma and taste — belong to olfaction and gustation: no photograph can
reveal a taste. And then there is an instructive case, rind thickness: it is physically visible,
yet the rind occupies a small portion of the slice, and we'll see that the model struggles to
"weight" it.
This taxonomy is not decorative: when we measure whether the model truly uses the image, we'll
find it succeeds exactly on the visible attributes and fails on the olfactory and gustatory ones.
It is a sanity check that validates the whole approach.
""")

# =====================================================================
# SLIDE 4 — TWO STEPS
# =====================================================================
s = newslide()
chrome(s, 4)
y0 = kicker_title(s, "The brief", "Two steps, one comparison goal")
card(s, Inches(0.7), y0+Inches(0.1), Inches(5.9), Inches(3.7),
     title="Step 1 · Text pre-processing", accent=PRIM,
     bullets=[("", "Clean and normalise telegraphic, dialectal, inconsistent descriptions"),
              ("", "Replace quantitative measures (mm/cm) with qualitative descriptions"),
              ("", "Produce captions as complete Italian sentences"),
              ("≈ 70%", " of the project's overall effort")],
     tsize=16.5, bsize=13.5)
card(s, Inches(6.75), y0+Inches(0.1), Inches(5.9), Inches(3.7),
     title="Step 2 · Three captioning methods", accent=GOLD,
     bullets=[("", "Apply three «as conceptually different as possible» encoder–decoders"),
              ("", "Focus on comparing methods, not on the single «best» model"),
              ("", "Multi-metric evaluation beyond BLEU alone"),
              ("≈ 30%", " of the project's overall effort")],
     tsize=16.5, bsize=13.5)
notes(s, """
The AI4FQC brief explicitly asks for two deliverables.
The first step is text pre-processing: take the tasters' descriptions — often telegraphic,
sometimes in dialect, occasionally inconsistent — and turn them into clean, normalised captions in
the form of Italian sentences. A specific, recurring requirement is to replace quantitative
measures, the rind's millimetres and centimetres, with qualitative descriptions. This step, as you
see, accounts for about 70% of the work: it is the most original part.
The second step is to apply and compare three encoder-decoder methods that are "as conceptually
different as possible". Note the emphasis: the brief does not ask for the best model, but for a
comparison of approaches. To do this honestly, we paired BLEU with a battery of complementary
metrics.
Consistently with these weights, I'll devote the first and larger part of the talk to building the
dataset, and the second to the models. Let's start with the data.
""")

# =====================================================================
# SLIDE 5 — SECTION PART I
# =====================================================================
s = newslide(DEEP)
rect(s, 0, 0, Inches(0.22), SH, GOLD)
cheese_holes(s, Inches(11.2), Inches(1.4), GOLD, base=0.7, alpha=18)
cheese_holes(s, Inches(2.0), Inches(6.2), GOLD, base=0.5, alpha=14)
_, tf = textbox(s, Inches(1.0), Inches(2.55), Inches(11), Inches(1.0))
setrun(tf.paragraphs[0], "PART I", 20, GOLD, bold=True, spacing=400)
_, tf = textbox(s, Inches(0.95), Inches(3.15), Inches(11.4), Inches(1.6))
setrun(tf.paragraphs[0], "Building the dataset", 46, CREAM, bold=True, font=FONT_H)
rect(s, Inches(1.0), Inches(4.55), Inches(2.4), Pt(4), GOLD)
_, tf = textbox(s, Inches(1.0), Inches(4.85), Inches(10.5), Inches(0.8))
setrun(tf.paragraphs[0],
       "From 51,988 raw, dialectal rows to 38,437 clean image–caption pairs",
       17, RGBColor(0xD9,0xD2,0xC2), italic=True)
notes(s, """
Let's enter the first part: building the dataset.
Here's the journey in one line: we start from nearly 52 thousand raw, heterogeneous and at times
dialectal rows, and we arrive at 38 thousand clean, normalised image-caption pairs in sentence
form. In between there's an eleven-stage pipeline and a precise strategy that let us do all this
for an AI cost of just 5 dollars and 60 cents.
This, in my view, is the most interesting part of the project from an engineering standpoint,
because the raw data was deceptively rich but deeply misaligned. Let's see why.
""")

# =====================================================================
# SLIDE 6 — RAW DATA & JOIN
# =====================================================================
s = newslide()
chrome(s, 6)
y0 = kicker_title(s, "The starting point", "Rich data, but misaligned")
card(s, Inches(0.7), y0, Inches(3.85), Inches(3.55),
     title="What we had", accent=PRIM, tsize=15,
     bullets=[("2,745", " BMP photos of sections (Slice / Grain)"),
              ("4", " Excel workbooks (2018–2021), one sheet per attribute"),
              ("1", " codebook: dairy ↔ product ↔ letter (16 dairies)")],
     bsize=12.5)
card(s, Inches(4.7), y0, Inches(3.85), Inches(3.55),
     title="The join problem", accent=RED, tsize=15,
     body="Images and comments «spoke different languages». The filename identifies the single "
          "wheel; the comment only the taster's tray. No row-level join — only a dairy-level join, "
          "via the codebook.",
     bsize=12.5)
card(s, Inches(8.7), y0, Inches(3.95), Inches(3.55),
     title="Invisible obstacles", accent=GOLD, tsize=15,
     bullets=[("", "Triple dairy indexing: TN_306 / TN306 / 306"),
              ("", "Inconsistent headers across years"),
              ("", "Italian decimal comma ('7,48')"),
              ("", "Provenance propagated downstream → traceability")],
     bsize=12)
notes(s, """
The starting material looked rich: nearly 2,750 high-resolution BMP photos of cheese sections, in
two views — the whole slice and a close-up of the grain; four Excel workbooks, one per year from
2018 to 2021, with one sheet per attribute; and a codebook mapping the sixteen dairies.
The problem, however, was not quantity but alignment. Images and comments, in practice, spoke
different languages. The photo filename identifies the single wheel — a precise sample ID — while
the taster's comment identifies only the tray, the "product". So there was no row-by-row join
between a photo and its comment: there was only a dairy-level join, going through the codebook.
On top of that, three concrete obstacles that cost real time: the same dairy appears written three
different ways; the Excel headers change from year to year; and the numeric scores use the Italian
decimal comma. Every provenance field is propagated downstream, so each final row stays traceable
back to the original comment.
""")

# =====================================================================
# SLIDE 7 — BROADCAST
# =====================================================================
s = newslide()
chrome(s, 7)
y0 = kicker_title(s, "The join solution", "Dairy-level broadcast")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.55))
setrun(tf.paragraphs[0],
       "Each comment is propagated to all coherent images of that dairy in that session: up to 4 rows from a single comment.",
       13.5, INK)
dy = y0+Inches(0.7)
rect(s, Inches(0.8), dy+Inches(0.9), Inches(2.7), Inches(1.3), RGBColor(0xF6,0xE3,0xDD), round_=True, shadow=True)
rect(s, Inches(0.8), dy+Inches(0.9), Inches(0.09), Inches(1.3), RED)
_, tf = textbox(s, Inches(1.0), dy+Inches(1.0), Inches(2.4), Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; setrun(p,"Panel comment",12,DEEP,bold=True)
p=tf.add_paragraph(); setrun(p,'«Sauerkraut»',15,RED,bold=True,italic=True,font=FONT_H)
p=tf.add_paragraph(); setrun(p,"(Smell · TN_306)",10.5,GREY)
labs=["P3a — Slice","P3a — Grain","P3b — Slice","P3b — Grain"]
for i,l in enumerate(labs):
    iy=dy+i*Inches(0.78)
    rect(s, Inches(5.0), iy, Inches(2.7), Inches(0.62), CREAM2, round_=True, shadow=True)
    rect(s, Inches(5.0), iy, Inches(0.08), Inches(0.62), PRIM)
    _, tf=textbox(s, Inches(5.2), iy, Inches(2.5), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0], l, 12.5, INK, bold=True)
    ar=s.shapes.add_connector(2, Inches(3.55), dy+Inches(1.55), Inches(4.95), iy+Inches(0.31))
    ar.line.color.rgb=GREY; ar.line.width=Pt(1.5)
rect(s, Inches(9.2), dy+Inches(0.9), Inches(3.3), Inches(1.3), RGBColor(0xE3,0xEE,0xE3), round_=True, shadow=True)
rect(s, Inches(9.2), dy+Inches(0.9), Inches(0.09), Inches(1.3), GREEN)
_, tf=textbox(s, Inches(9.4), dy+Inches(0.95), Inches(3.0), Inches(1.2), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; setrun(p,"4 image–caption rows",13,DEEP,bold=True)
p=tf.add_paragraph(); setrun(p,"(same caption)",11.5,GREEN)
for i in range(4):
    iy=dy+i*Inches(0.78)
    ar=s.shapes.add_connector(2, Inches(7.75), iy+Inches(0.31), Inches(9.15), dy+Inches(1.55))
    ar.line.color.rgb=GREY; ar.line.width=Pt(1.5)
rect(s, Inches(0.8), dy+Inches(2.55), Inches(11.7), Inches(0.7), DEEP, round_=True)
_, tf=textbox(s, Inches(1.0), dy+Inches(2.55), Inches(11.3), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]
setrun(p,"Net effect:  ",12.5,GOLD,bold=True)
setrun(p,"51,988 rows (image × panelist × attribute) · 39,510 with a comment · pairing rate ≈ 76%",12.5,CREAM)
notes(s, """
How do we solve the lack of a wheel-level join? With what we called a dairy-level broadcast. The
idea is to propagate each taster's comment to all coherent images of that dairy in that session.
Look at the example. A single comment about smell — the word "Sauerkraut" for dairy TN_306 —
becomes the caption of four different images: replicas a and b, each in the two views, slice and
grain. From one comment we get up to four image-caption rows, with the same caption.
One might object: isn't this artificial duplication? Actually no — for training it's an advantage:
more examples of the same visual-textual association make the signal more robust.
The net effect of the join is nearly 52 thousand rows, obtained as image times panelist times
attribute, of which about 39,500 have a non-empty comment: a pairing rate around 76%, lower in the
2021 sessions where some fields were left blank. From here the pipeline begins.
""")

# =====================================================================
# SLIDE 8 — PIPELINE 11 STAGES
# =====================================================================
s = newslide()
chrome(s, 8)
y0 = kicker_title(s, "The architecture", "An eleven-stage pipeline")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.5))
p=tf.paragraphs[0]
setrun(p,"Guiding principle: ",13.5,INK,bold=True)
setrun(p,"deterministic-before-the-LLM — everything that can be done with reproducible, free code is done before calling the model.",13.5,INK)
phases=[("0","Unified table","51,988",PRIM),("1","Deterministic prep","39,356",PRIM),
        ("2–3","Vocabulary + audit","7 attr.",PRIM),("4","Cleanup + qualitative","7,705 unique",PRIM),
        ("5","Noise drop","7,689",PRIM),("6–7","Prompt + pilot","105",GOLD_D),
        ("8","LLM batch","7,689",GOLD_D),("9","Manual salvage","+916",GREEN),
        ("10","Broadcast + sentence","38,437",GREEN)]
cols=3; cw=Inches(3.9); ch=Inches(1.18); gx=Inches(0.15); gy=Inches(0.18)
x0=Inches(0.75); yA=y0+Inches(0.65)
for i,(n,t,v,col) in enumerate(phases):
    r=i//cols; c=i%cols
    x=x0+c*(cw+gx); y=yA+r*(ch+gy)
    rect(s, x, y, cw, ch, CREAM2, round_=True, shadow=True)
    oval(s, x+Inches(0.18), y+Inches(0.3), Inches(0.58), Inches(0.58), col)
    _, tf=textbox(s, x+Inches(0.18), y+Inches(0.3), Inches(0.58), Inches(0.58), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,n,12.5,WHITE,bold=True)
    _, tf=textbox(s, x+Inches(0.92), y+Inches(0.16), cw-Inches(1.05), ch-Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; setrun(p,t,13.5,DEEP,bold=True,font=FONT_H)
    p=tf.add_paragraph(); setrun(p,v,12,col if col!=GOLD_D else GOLD_D,bold=True)
_, tf=textbox(s, Inches(0.75), yA+3*(ch+gy)-Inches(0.02), Inches(11.8), Inches(0.4))
p=tf.paragraphs[0]
setrun(p,"■ ",13,PRIM,bold=True); setrun(p,"deterministic (free)    ",11.5,INK)
setrun(p,"■ ",13,GOLD_D,bold=True); setrun(p,"with LLM    ",11.5,INK)
setrun(p,"■ ",13,GREEN,bold=True); setrun(p,"final assembly",11.5,INK)
notes(s, """
Here's the full map: eleven stages, numbered zero to ten, each with its own Python script, explicit
inputs and outputs, and an inspectable report.
The guiding principle — and this is the most important architectural decision of the first part —
is "deterministic before the LLM". It means: everything that can be done with reproducible,
verifiable, free code, we do before calling the language model.
In blue you see the deterministic stages: from the unified table, to text preparation, to building
the vocabulary, to cleanup and deduplication. In gold the stages that involve the LLM: prompt
design, the pilot and the full batch. In green the final assembly.
Follow the numbers along the stages: 52 thousand initial rows, reduced and normalised, then
compressed to about 7,700 unique captions — and only this subset reaches the paid model. It is
exactly this compression that keeps the cost to a few dollars. Let's look at the key stages one by
one.
""")

# =====================================================================
# SLIDE 9 — FUNNEL (figure)
# =====================================================================
s = newslide()
chrome(s, 9)
y0 = kicker_title(s, "The data flow", "The funnel: where rows are lost (and kept)")
pic_fit(s, A("fig_funnel.png"), Inches(0.7), y0, Inches(7.7), Inches(4.4))
card(s, Inches(8.7), y0+Inches(0.1), Inches(3.95), Inches(4.1),
     title="Reading", accent=GOLD, tsize=15,
     bullets=[("−12,632", " rows in deterministic prep"),
              ("12,478", " empty / null"),
              ("86", " meta-comments  ·  68 near-empty"),
              ("75.7%", " retention after stage 1")],
     bsize=12.5)
notes(s, """
This is the same story, seen as a funnel. The axis shows the number of training rows as they move
through the pipeline.
The biggest drop — over 12,600 rows — happens entirely in the deterministic preparation of stage 1.
But watch what we discard: the vast majority, nearly 12,500 rows, are empty or null values. Only 86
are meta-comments and 68 are near-empty fragments. In other words, we are not discarding sensory
content: we are removing certain noise.
The retention rate after stage 1 is 75.7%. And here's a deliberate choice I want to stress: at this
stage, drops must be certain. A sentence like "I don't penalise it, but it smells of barn" contains
a real descriptor and must NOT be discarded by a regex — that kind of ambiguity is better handled
by the LLM later. The deterministic blacklist is therefore small and conservative by design.
""")

# =====================================================================
# SLIDE 10 — STAGE 1
# =====================================================================
s = newslide()
chrome(s, 10)
y0 = kicker_title(s, "Stage 1", "Deterministic prep: clean without rewriting")
card(s, Inches(0.7), y0, Inches(5.85), Inches(3.6),
     title="Operations", accent=PRIM, tsize=15.5,
     bullets=[("", "Filter empties/N-A + Unicode normalisation (NFC)"),
              ("", "Remove non-breaking spaces, zero-width, tabs, newlines"),
              ("", "Drop meta-comments (small, conservative blacklist)"),
              ("", "Drop near-empty noise (< 2 alphanumeric characters)"),
              ("", "Dual column raw / norm → reversibility and audit")],
     bsize=13)
card(s, Inches(6.7), y0, Inches(5.95), Inches(3.6),
     title="Philosophy", accent=GOLD, tsize=15.5, fill=DEEP, tcolor=GOLD,
     body="«Drops must be certain». The LLM will handle ambiguity better: a negation with a "
          "descriptor — «I don't penalise it, but it smells of barn» — carries real information "
          "and is kept.\n\nResult: 39,356 rows (75.7% retention), 12,632 drops, of which 12,478 "
          "empty, 86 meta, 68 near-empty.",
     bsize=13)
notes(s, """
Let's enter stage 1, the deterministic preparation. The golden rule here is: clean without
rewriting. We don't change the taster's words, we only clean them.
The operations are all mechanical and verifiable: we filter empties and N/A, normalise Unicode to
NFC form, remove non-breaking spaces, zero-width characters, tabs and newlines. Then we remove
meta-comments with a deliberately small blacklist, and near-empty fragments, those with fewer than
two alphanumeric characters after cleaning. We always keep two columns, raw text and normalised
text, so every transformation is reversible and auditable.
The philosophy, on the right, is the key: drops must be certain. Anything ambiguous passes through,
because the LLM will handle it better than any fixed rule. The net result is 39,356 rows, with the
drops that — I repeat because it matters — are 99% simply empties.
""")

# =====================================================================
# SLIDE 11 — STAGES 2-3 VOCAB (figure)
# =====================================================================
s = newslide()
chrome(s, 11)
y0 = kicker_title(s, "Stages 2–3", "A controlled vocabulary, tailor-made")
pic_fit(s, A("fig_vocab.png"), Inches(0.7), y0, Inches(7.7), Inches(4.4))
card(s, Inches(8.7), y0+Inches(0.1), Inches(3.95), Inches(4.1),
     title="Why custom", accent=PRIM, tsize=15,
     bullets=[("", "Italian NLP tools err: panna→panno, latte→latto"),
              ("", "Useful stopwords here: «molto», «poco», «leggermente»"),
              ("", "Deterministic lemmatiser + sing./plur. merge"),
              ("Dual use:", " audit basis + stylistic anchor in the prompt")],
     bsize=12)
notes(s, """
Stages 2 and 3 build a controlled vocabulary for each of the seven attributes: the most frequent
lemmas and bigrams, with a tailor-made Italian lemmatisation.
Why not use a standard NLP library? Because on the sensory lexicon they err systematically: they
turn "panna" (cream) into "panno" (cloth), "latte" (milk) into "latto". And above all they include
among the stopwords words that here are informative: "molto", "poco", "leggermente" — very, little,
slightly — are sensory intensifiers, not noise. So we wrote a deterministic lemmatiser, with maps
of abbreviations and typos and a singular-plural merge that fires only when both forms are attested
in the corpus.
In the chart you see the vocabulary breadth per attribute: paste structure is the richest, with
nearly 800 lemmas and over 40 thousand tokens; rind thickness is the poorest, consistent with its
lower descriptive variety. The vocabulary serves two purposes: as a basis for the audit, and as a
stylistic anchor in the LLM prompt — not as a closed dictionary, since the model already speaks
Italian.
""")

# =====================================================================
# SLIDE 12 — STAGE 4 NUMBERS->QUALITY
# =====================================================================
s = newslide()
chrome(s, 12)
y0 = kicker_title(s, "Stage 4 · key requirement", "From numbers to quality: the rind")
_, tf = textbox(s, Inches(0.7), y0-Inches(0.05), Inches(11.9), Inches(0.5))
setrun(tf.paragraphs[0],
       "Hundreds of captions were just measures («10», «1 cm», «8-10 mm»). A deterministic function converts everything to mm and assigns a bucket.",
       13.5, INK)
rows=[("Threshold (mm)","Bucket"),("< 8","Very thin"),("8 ≤ x < 10","Thin"),
      ("10 ≤ x < 14","Medium"),("14 ≤ x < 18","Thick"),("≥ 18","Very thick")]
tx=Inches(0.7); ty=y0+Inches(0.7); tw=Inches(4.7); th=Inches(3.4)
tbl=s.shapes.add_table(len(rows),2,tx,ty,tw,th).table
tbl.columns[0].width=Inches(2.2); tbl.columns[1].width=Inches(2.5)
for ri,(a,b) in enumerate(rows):
    for ci,val in enumerate((a,b)):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ri==0:
            setrun(p,val,13,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            setrun(p,val,12.5,INK,bold=(ci==1))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
rect(s, Inches(5.75), ty-Inches(0.05), Inches(6.9), Inches(3.6), WHITE, round_=True, shadow=True)
pic_fit(s, A("fig_spessore.png"), Inches(5.95), ty+Inches(0.1), Inches(6.5), Inches(3.3))
notes(s, """
Stage 4 tackles head-on the brief's explicit requirement: replacing quantitative descriptions with
qualitative ones. The emblematic case is rind thickness, where hundreds of captions were plain
numbers: "10", "1 cm", "on average 9 mm", "8-10 mm".
The solution is a deterministic function that recognises fully numeric captions, converts
everything to millimetres — with a heuristic reading values below 5 as centimetres — and assigns a
qualitative bucket per the table you see: below 8 millimetres "very thin", up to 10 "thin", and so
on up to "very thick" beyond 18.
Why do it deterministically and not with the LLM? Because it eliminates at the root the model's main
source of inconsistency, which otherwise bucketed "1 cm" and "10 mm" — the exact same physical
measure — differently. In the chart on the right you see the result: 424 numeric rows neatly
collapsed into five buckets. Mixed captions, those blending numbers and description, are instead
left to the LLM because they require contextual reasoning.
""")

# =====================================================================
# SLIDE 13 — DEDUP (figure)
# =====================================================================
s = newslide()
chrome(s, 13)
y0 = kicker_title(s, "Stage 4 · the decisive trick", "Deduplication: 5.1× compression")
pic_fit(s, A("fig_dedup.png"), Inches(0.7), y0, Inches(7.7), Inches(4.4))
card(s, Inches(8.7), y0+Inches(0.1), Inches(3.95), Inches(4.1),
     title="Why it matters", accent=GREEN, tsize=15,
     bullets=[("", "The broadcast replicates the same text many times"),
              ("39,356", " rows → 7,705 unique captions"),
              ("5.1×", " compression by (caption, attribute) key"),
              ("", "Dominant 4× pattern: 6,595 captions exactly 4 times"),
              ("→", " cuts the LLM cost ~5×")],
     bsize=12)
notes(s, """
This is the single most important trick for the project's economics: deduplication.
Remember the broadcast? By replicating each comment across four images, the exact same text would
have been sent to the LLM many times, paying for it each time. That makes no sense. So, before
calling the model, we deduplicate by key: the caption-plus-attribute pair.
The result is sharp: over 39 thousand rows compress into 7,705 unique captions, a 5.1-fold
compression. The chart shows, per attribute, prepared rows versus unique captions. The dominant
pattern is exactly the broadcast's factor of four: nearly 6,600 captions appear precisely four
times.
The practical consequence is enormous: we send the LLM only the unique captions, cutting the cost
by about five times. It's the perfect illustration of the "deterministic before the LLM" principle:
a simple deduplication, for free, cuts the paid component's cost by a factor of five.
""")

# =====================================================================
# SLIDE 14 — STAGES 6-8 LLM
# =====================================================================
s = newslide()
chrome(s, 14)
y0 = kicker_title(s, "Stages 6–8", "LLM rewriting: the genuinely hard work")
card(s, Inches(0.7), y0, Inches(3.85), Inches(3.6),
     title="6 · Prompt design", accent=GOLD, tsize=15,
     bullets=[("", "Per-attribute system prompt (~5 KB)"),
              ("11", " rules (zero invention, NON_DESCRITTO escape)"),
              ("Top-60", " lemmas as a style anchor"),
              ("6", " real few-shot examples")], bsize=12)
card(s, Inches(4.7), y0, Inches(3.85), Inches(3.6),
     title="7 · Pilot run", accent=PRIM, tsize=15,
     bullets=[("105", " captions (15 per attribute)"),
              ("", "Found 2 failure modes: 1cm vs 10mm; empty format"),
              ("", "Both fixed in the prompt"),
              ("3 workers", " = optimal; 8 → storms of 429s")], bsize=12)
card(s, Inches(8.7), y0, Inches(3.95), Inches(3.6),
     title="8 · Full batch", accent=GREEN, tsize=15,
     bullets=[("7,689 / 7,689", " completed, 0 errors"),
              ("~25–30", " minutes"),
              ("Haiku 4.5", " · Anthropic Batch API"),
              ("~$4.50", " for the batch")], bsize=12)
notes(s, """
We reach the "smart" core of the pipeline: LLM rewriting, stages 6, 7 and 8. Here the LLM does the
work no regex could do faithfully: expanding "Sauerkraut" into "Smell of sauerkraut", normalising
dialect, removing meta-comments while preserving the descriptor, turning questions into statements.
Stage 6 is prompt design: a system prompt for each attribute, about 5 kilobytes, with eleven rules
— including a "zero invention" rule and an escape, NON_DESCRITTO, for contentless inputs — the
sixty most frequent lemmas as a stylistic anchor, and six real few-shot examples.
Stage 7, the pilot on 105 captions, served to discover two failure modes before spending: the
centimetre-millimetre inconsistency and some format violations. Both fixed in the prompt. A concrete
operational lesson: three parallel workers were optimal, eight triggered rate-limit storms.
Stage 8, the full batch: all 7,689 unique captions sent as a single job. Result: 7,689 out of 7,689
completed, zero errors, in half an hour, for about four and a half dollars. Which brings us to the
question: why Haiku?
""")

# =====================================================================
# SLIDE 15 — COST (native chart)
# =====================================================================
s = newslide()
chrome(s, 15)
y0 = kicker_title(s, "Model choice", "Why Haiku 4.5: same job, a fraction of the cost")
cd = CategoryChartData()
cd.categories = ["Haiku 4.5", "Sonnet", "Opus"]
cd.add_series("Estimated cost ($)", (5.60, 13.50, 67.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.7), y0+Inches(0.05), Inches(7.4), Inches(4.3), cd)
ch = gf.chart; style_chart(ch, legend=False)
plot=ch.plots[0]; plot.has_data_labels=True
plot.data_labels.number_format='"$"0.00'; plot.data_labels.number_format_is_linked=False
plot.data_labels.font.size=Pt(13); plot.data_labels.font.bold=True; plot.data_labels.font.color.rgb=DEEP
plot.data_labels.position=XL_LABEL_POSITION.OUTSIDE_END
plot.gap_width=80
points_color(plot.series[0], [GREEN, GOLD, RED])
ch.value_axis.has_major_gridlines=True
card(s, Inches(8.35), y0+Inches(0.1), Inches(4.3), Inches(4.15),
     title="The pilot's verdict", accent=GREEN, tsize=15.5,
     bullets=[("2.4×", " cheaper than Sonnet"),
              ("12×", " cheaper than Opus"),
              ("", "On the verified task, bigger models gave no measurable extra quality"),
              ("$5.60", " total LLM cost of the project")], bsize=13)
notes(s, """
One of the project's most pragmatic decisions. For a well-bounded rewriting task, already verified
on the pilot, which is the right model?
The chart answers. The same work would cost about 5 dollars 60 with Haiku 4.5, versus about 13.50
with Sonnet and over 67 with Opus. Haiku is 2.4 times cheaper than Sonnet and twelve times cheaper
than Opus.
The real question, though, is not "which is cheapest" but "is the cheap model good enough?". This is
where the pilot comes in: having tested it on a stratified sample, we verified that for this narrow
task the more expensive models offered no measurable quality gain. Spending more would have been
throwing money away.
The total LLM cost of the entire project — pilot plus batch — was 5 dollars and 60 cents. It is
concrete proof that good upstream data engineering is worth more than an expensive downstream model.
""")

# =====================================================================
# SLIDE 16 — OUTPUT VALIDATION
# =====================================================================
s = newslide()
chrome(s, 16)
y0 = kicker_title(s, "Programmatic validation", "Step 1 is satisfied 100%")
rows=[("Check","Violations"),
      ("Output starts with the expected attribute prefix","1*"),
      ("Output contains digits","0"),
      ("Output contains units (mm / cm / %)","0"),
      ("Output longer than 25 words","0"),
      ("Output empty","0"),
      ("Output multi-line or with markup","0")]
tx=Inches(0.7); ty=y0+Inches(0.05); tw=Inches(7.7); th=Inches(3.9)
tbl=s.shapes.add_table(len(rows),2,tx,ty,tw,th).table
tbl.columns[0].width=Inches(6.0); tbl.columns[1].width=Inches(1.7)
for ri,(a,b) in enumerate(rows):
    for ci,val in enumerate((a,b)):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ci==1: p.alignment=PP_ALIGN.CENTER
        if ri==0:
            setrun(p,val,13.5,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            col = GREEN if val=="0" else GOLD_D
            setrun(p,val,13 if ci==0 else 16, INK if ci==0 else col, bold=(ci==1))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(10)
card(s, Inches(8.7), y0+Inches(0.05), Inches(3.95), Inches(3.95),
     title="Zero numbers,\nzero units", accent=GREEN, tsize=18, fill=DEEP, tcolor=GREEN,
     body="\nThe quantitative → qualitative conversion succeeded 100%: the main requirement of "
          "Step 1 is fully met.\n\n* the only «violation» is an alternative, still-valid prefix.", bsize=13)
notes(s, """
How do we know Step 1 is truly satisfied? Not by impression, but with programmatic validation on all
7,689 model outputs.
The table summarises the checks. The output always starts with the expected attribute prefix — with
a single "violation" that is actually a valid alternative prefix. Above all, look at the numbers
that matter for the brief: zero outputs with digits, zero outputs with units like millimetres or
centimetres, zero overly long outputs, zero empty outputs, zero outputs with markup. All zeros.
This means the conversion from quantitative to qualitative — the core requirement of the first step
— succeeded one hundred percent: in the final captions there remains no number and no unit of
measure. It's a clean, verifiable fulfilment of the deliverable, not a promise. And when a reviewer
asks "but are you sure you removed all the measures?", the answer is a table, not an opinion.
""")

# =====================================================================
# SLIDE 17 — SALVAGE (figure)
# =====================================================================
s = newslide()
chrome(s, 17)
y0 = kicker_title(s, "Stage 9", "Manual salvage: cleaning ≠ throwing away")
pic_fit(s, A("fig_salvage.png"), Inches(0.7), y0, Inches(7.4), Inches(4.4))
card(s, Inches(8.35), y0+Inches(0.1), Inches(4.3), Inches(4.1),
     title="Hand-curation", accent=GOLD, tsize=15,
     bullets=[("", "291/362 NON_DESCRITTO had a vocabulary lemma → LLM too cautious"),
              ("178", " captions salvaged by hand"),
              ("«marcio, putrido,»", " → «Smell of rot and putridity.»"),
              ("362 → 184", " unique NON_DESCRITTO"),
              ("+916", " training rows recovered")], bsize=12)
notes(s, """
Stage 9 is one of those that separate a curated dataset from a merely "cleaned" one. After the
batch, a scan revealed something interesting: 291 of the 362 captions the LLM had labelled
NON_DESCRITTO — i.e. "there's no description here" — actually contained at least one lemma from the
controlled vocabulary. A sign the model had been overly cautious with the escape rule on borderline
inputs, those mixing a judgement and a descriptor.
So we hand-curated a salvage map of 178 captions where the descriptor was real and faithful.
Examples: "marcio, putrido," becomes "Smell of rot and putridity"; "Sangue,,," becomes "Aroma of
blood".
The effect, in the chart: unique NON_DESCRITTO drop from 362 to 184, and this recovers 916 training
rows we would otherwise have lost. The message is simple: cleaning data does not mean throwing it
away. A little manual care, on a targeted subset, cost less than another LLM round and returned
nearly a thousand valid examples.
""")

# =====================================================================
# SLIDE 18 — FINAL DATASET (figure + schema)
# =====================================================================
s = newslide()
chrome(s, 18)
y0 = kicker_title(s, "Stage 10 · the deliverable", "The final dataset")
pic_fit(s, A("fig_final_rows.png"), Inches(0.7), y0, Inches(6.7), Inches(3.5))
kpis=[("38,437","final rows",PRIM),("1,497","unique images",GOLD_D),("2","caption forms",GREEN)]
for i,(v,l,c) in enumerate(kpis):
    x=Inches(0.7)+i*Inches(2.3)
    rect(s, x, y0+Inches(3.65), Inches(2.15), Inches(0.95), CREAM2, round_=True, shadow=True)
    _, tf=textbox(s, x, y0+Inches(3.7), Inches(2.15), Inches(0.9), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,v,19,c,bold=True,font=FONT_H)
    p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; setrun(p,l,11,INK)
card(s, Inches(7.6), y0, Inches(5.05), Inches(4.6),
     title="Delivered schema", accent=PRIM, tsize=15,
     bullets=[("captions_final.csv", " — 38,437 × 18 col. (provenance)"),
              ("image_caption_attribute.csv", " — 4-column view"),
              ("by_attribute/<Attr>.csv", " — 7 per-attribute splits"),
              ("caption", " (compact, ~4–8 words)"),
              ("caption_sentence", " (Italian sentence, ~7–15 words)"),
              ("", "Deterministic sentence form: 100% template, 0 extra LLM round-trips")], bsize=11.5)
notes(s, """
Here we are at the first part's deliverable. Stage 10 reconnects the clean captions to the broadcast
table and produces the final dataset: 38,437 training rows over 1,497 unique images.
The chart shows the row distribution per attribute, and note there's significant variance — paste
structure has nearly double the rows of rind thickness. This imbalance will matter in the
cross-attribute comparison of the second part.
On the format side, we made a flexibility choice: each row carries two caption forms. A compact one,
four to eight words, anchored to the attribute, like "Smell of cream". And a complete declarative
Italian sentence, "The cheese has a smell of cream", better suited to encoder-decoders and to
captioning metrics. This sentence transformation is fully deterministic, with no further LLM round.
The dataset is delivered in several forms — full table with all provenance, simplified four-column
view, and seven per-attribute splits — to serve any downstream architecture. With this we close the
data part and move on to the models.
""")

# =====================================================================
# SLIDE 19 — SECTION PART II
# =====================================================================
s = newslide(DEEP)
rect(s, 0, 0, Inches(0.22), SH, GOLD)
cheese_holes(s, Inches(11.0), Inches(5.6), GOLD, base=0.7, alpha=18)
cheese_holes(s, Inches(2.2), Inches(1.4), GOLD, base=0.5, alpha=14)
_, tf = textbox(s, Inches(1.0), Inches(2.55), Inches(11), Inches(1.0))
setrun(tf.paragraphs[0], "PART II", 20, GOLD, bold=True, spacing=400)
_, tf = textbox(s, Inches(0.95), Inches(3.15), Inches(11.6), Inches(1.6))
setrun(tf.paragraphs[0], "Model selection and training", 42, CREAM, bold=True, font=FONT_H)
rect(s, Inches(1.0), Inches(4.5), Inches(2.4), Pt(4), GOLD)
_, tf = textbox(s, Inches(1.0), Inches(4.8), Inches(10.8), Inches(0.8))
setrun(tf.paragraphs[0],
       "Three encoder–decoder methods, one determinant: the visual encoder",
       17, RGBColor(0xD9,0xD2,0xC2), italic=True)
notes(s, """
Second part: the models.
Here the brief's goal is the comparison of methods, and let me give you the conclusion upfront, so
you can follow the data with the right key. Among everything we could change — the decoder, the
amount of data, the hyperparameters — the single factor that decides whether the model truly uses
the image is one: the visual encoder.
We'll build this conclusion step by step, and verify it with two independent tools: the shuffle test
and the CLIPScore. Let's start with the three architectures.
""")

# =====================================================================
# SLIDE 20 — THE THREE METHODS
# =====================================================================
s = newslide()
chrome(s, 20)
y0 = kicker_title(s, "The architectures", "Three methods, two axes of variation")
models=[("m1","ResNet-50\n(CNN, frozen)","LSTM\nfrom scratch",M1C),
        ("m3","ViT-B/16\n(frozen)","Transformer\nfrom scratch",M3C),
        ("m6","ViT-B/16\n(frozen)","GePpeTto\npretrained Italian GPT-2",M6C)]
cw=Inches(3.6); gx=Inches(0.55); x0=Inches(0.95); yT=y0+Inches(0.25)
for i,(m,enc,dec,col) in enumerate(models):
    x=x0+i*(cw+gx)
    _, tf=textbox(s, x, yT, cw, Inches(0.5)); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    setrun(p,m,22,col,bold=True,font=FONT_H)
    rect(s, x, yT+Inches(0.55), cw, Inches(1.15), CREAM2, round_=True, shadow=True)
    rect(s, x, yT+Inches(0.55), cw, Inches(0.1), PRIM)
    _, tf=textbox(s, x+Inches(0.15), yT+Inches(0.6), cw-Inches(0.3), Inches(1.05), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,"ENCODER",10,PRIM,bold=True,spacing=150)
    for ln in enc.split("\n"):
        p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; setrun(p,ln,13,DEEP,bold=True)
    a=s.shapes.add_connector(2, x+cw/2, yT+Inches(1.75), x+cw/2, yT+Inches(2.05))
    a.line.color.rgb=GREY; a.line.width=Pt(2)
    rect(s, x, yT+Inches(2.1), cw, Inches(1.15), CREAM2, round_=True, shadow=True)
    rect(s, x, yT+Inches(2.1), cw, Inches(0.1), GOLD)
    _, tf=textbox(s, x+Inches(0.15), yT+Inches(2.15), cw-Inches(0.3), Inches(1.05), anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p,"DECODER",10,GOLD_D,bold=True,spacing=150)
    for ln in dec.split("\n"):
        p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; setrun(p,ln,12.5,DEEP,bold=True)
_, tf=textbox(s, x0, yT+Inches(3.45), Inches(8.1), Inches(0.4))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"↔ m1 vs m3 isolates the ENCODER        ↔ m3 vs m6 isolates the DECODER",12.5,RED,bold=True)
_, tf=textbox(s, Inches(0.95), yT+Inches(3.9), Inches(11.5), Inches(0.5))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"All encoders are frozen: only decoder and cross-attention projection are trained → fair comparison.",12,INK,italic=True)
notes(s, """
The brief asks for three methods "as conceptually different as possible". We chose them to isolate
the three main axes of variation in captioning.
m1 is the classic: a convolutional encoder, ResNet-50, with an LSTM decoder trained from scratch.
m3 changes the encoder: a Vision Transformer, ViT-B/16, with a Transformer decoder, still from
scratch. m6 changes the decoder: the same ViT, but the decoder is GePpeTto, an already-pretrained
Italian GPT-2.
The beauty of this design is in the axes, at the bottom. Comparing m1 with m3 isolates the encoder's
effect, because the decoder stays the same "from scratch" type. Comparing m3 with m6 isolates the
decoder's effect, because the encoder is identical. It's a small controlled experiment.
A crucial choice for fairness: all encoders are frozen. We train only the decoder and the
cross-attention projection. So each decoder receives exactly the same visual features, and the
differences we observe are attributable to architecture, not to a lucky encoder. This, we'll see,
is also the main limitation — and the main pointer to future work.
""")

# =====================================================================
# SLIDE 21 — SETUP
# =====================================================================
s = newslide()
chrome(s, 21)
y0 = kicker_title(s, "Experimental setup", "Training on Kaggle, fair comparison")
rows=[("Model","epochs","batch","lr","scheduler","patience"),
      ("m1 · CNN+LSTM","50","32","3e-4","StepLR","7"),
      ("m3 · ViT+Transf.","30","16","1e-4","cosine","5"),
      ("m6 · ViT+GePpeTto","20","8","5e-5","cosine","5")]
tx=Inches(0.7); ty=y0; tw=Inches(7.6); th=Inches(2.1)
tbl=s.shapes.add_table(len(rows),6,tx,ty,tw,th).table
widths=[2.3,1.0,1.0,1.1,1.2,1.0]
for ci,w in enumerate(widths): tbl.columns[ci].width=Inches(w)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ci>0: p.alignment=PP_ALIGN.CENTER
        if ri==0:
            setrun(p,val,11.5,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            setrun(p,val,11.5,INK,bold=(ci==0))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(6)
card(s, Inches(0.7), ty+Inches(2.35), Inches(7.6), Inches(1.95),
     title="Fair-evaluation choices", accent=PRIM, tsize=14.5,
     bullets=[("Sample-disjoint split:", " 674 train / 143 val / 147 test — no leak"),
              ("Nucleus sampling", " (top-p 0.9, T 0.7): beam collapses to the modal caption"),
              ("", "Healthy curves: no overfitting, NaN or collapse")], bsize=12)
card(s, Inches(8.45), y0, Inches(4.2), Inches(4.3),
     title="Kaggle lesson", accent=GOLD, tsize=15, fill=DEEP, tcolor=GOLD,
     body="The enable_gpu flag in the metadata does NOT actually enable the GPU: it must be turned "
          "on manually.\n\nThe free tier tends to the P100 (sm_60): the preinstalled PyTorch must be "
          "reinstalled from a cu118 wheel.\n\nAll training: free T4 / P100 GPU.", bsize=12.5)
notes(s, """
The experimental setup, briefly. All training ran on Kaggle's free GPU tier, with T4 or P100. The
hyperparameters in the table are reasonable defaults for each architecture, not the fruit of an
exhaustive search: m1 more epochs and a higher learning rate, m6 few epochs and a small learning
rate, as is natural for a pretrained model that only needs fine-tuning.
Two choices ensure a fair comparison. First, the split is sample-disjoint: the same cheese never
appears both in training and in test, so we avoid leakage. Second, at inference we use nucleus
sampling instead of beam search, because on small datasets beam collapses to the most frequent
caption, artificially inflating BLEU but reducing diversity.
The training curves are healthy across all three models — no overfitting, no NaN, no collapse — so
the differences in results are not due to training pathologies.
On the right, a Kaggle lesson learned the hard way, for anyone wanting to reproduce: the enable_gpu
flag in the metadata is not enough, and the preinstalled PyTorch on P100s must be reinstalled.
Details, but the kind that cost you an afternoon.
""")

# =====================================================================
# SLIDE 22 — BLEU PER ATTRIBUTE (figure)
# =====================================================================
s = newslide()
chrome(s, 22)
y0 = kicker_title(s, "Results · per-attribute", "BLEU-4: models win where captions are diverse")
pic_fit(s, A("fig_bleu_attr.png"), Inches(0.7), y0, Inches(8.1), Inches(4.4))
card(s, Inches(9.0), y0+Inches(0.1), Inches(3.65), Inches(4.1),
     title="Reading", accent=PRIM, tsize=15,
     bullets=[("Win:", " Structure (+0.128), Smell, Texture, Colour"),
              ("Tie/lose:", " Aroma, Taste, Thickness"),
              ("", "Not from weakness: BLEU rewards the modal predictor on low-diversity captions"),
              ("4 / 7", " attributes beat the constant baseline")], bsize=12)
notes(s, """
Let's see the first results, in the per-attribute regime: seven separate models, one per sensory
dimension. The chart shows BLEU-4 per attribute, compared with the most_frequent baseline — the one
that always emits the most frequent caption.
The picture is sharp and two-sided. The trained models clearly beat the constant baseline on four
attributes: paste structure, with a notable margin of over 0.12 points, then smell, texture and
colour. On the other three — aroma, taste and rind thickness — they tie or lose.
Be careful not to misread this "lose". It does not mean the models are poor on those attributes. It
means that there the reference captions are not very diverse, very template-heavy, and under those
conditions BLEU rewards whoever always predicts the majority. It's a flaw of the metric, not of the
model — and we'll prove it shortly with the shuffle test, which on thickness confirms m3 and m6 do
use the image.
In short: 4 of 7 attributes beat the baseline, and the 3 "lost" ones are a BLEU artifact.
""")

# =====================================================================
# SLIDE 23 — GLOBAL (native chart)
# =====================================================================
s = newslide()
chrome(s, 23)
y0 = kicker_title(s, "Results · global model", "A single model across all attributes")
cd = CategoryChartData()
cd.categories = ["BLEU-1","BLEU-4","METEOR","ROUGE-L","CIDEr"]
cd.add_series("m1 · CNN+LSTM", (0.3501,0.1283,0.2938,0.2951,0.2517))
cd.add_series("m3 · ViT+Transf.", (0.3649,0.1237,0.2875,0.2981,0.2856))
cd.add_series("m6 · ViT+GePpeTto", (0.3657,0.1307,0.2928,0.3007,0.2741))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), y0, Inches(8.3), Inches(4.4), cd)
ch=gf.chart; style_chart(ch); color_series(ch.plots[0],[M1C,M3C,M6C])
ch.plots[0].gap_width=90; ch.value_axis.has_major_gridlines=True
ch.value_axis.maximum_scale=0.4
card(s, Inches(9.15), y0+Inches(0.1), Inches(3.5), Inches(4.1),
     title="m6 wins narrowly", accent=M6C, tsize=15,
     bullets=[("", "m6 first on BLEU-1/4, ROUGE-L, METEOR≈"),
              ("", "m3 first on CIDEr"),
              ("0.007", " BLEU-4 spread: a very tight cluster"),
              ("⚠", " most_frequent has the highest BLEU-1 but worst BLEU-4: degenerate")], bsize=12)
notes(s, """
Let's move to the global regime: a single model trained across all attributes together. It's an
intrinsically harder task, because the model must also choose which sensory dimension to describe,
and it works with a much larger vocabulary.
The chart compares the three models on five metrics. Two observations. First, m6 — the ViT with
GePpeTto — wins narrowly almost everywhere: it's first on BLEU-1, BLEU-4 and ROUGE-L, and
essentially tied on METEOR; m3 only beats it on CIDEr. Second, and equally important, the three
models are in a very tight cluster: the BLEU-4 spread is just seven thousandths. They are very close.
A note on the baseline, which I left out of the chart for readability but is crucial: most_frequent,
the constant sentence, would have the highest BLEU-1 of all but the worst BLEU-4. It's the
degenerate behaviour of "always predict the majority": you guess many single words, but you never
compose a correct four-word sequence different from the template. Keep this anomaly in mind: it's
the first clue that BLEU, on its own, is misleading us.
""")

# =====================================================================
# SLIDE 24 — PER-ATTR VS GLOBAL (figure)
# =====================================================================
s = newslide()
chrome(s, 24)
y0 = kicker_title(s, "Mind the interpretation", "Per-attribute vs global: two different rulers")
pic_fit(s, A("fig_perattr_vs_global.png"), Inches(0.7), y0, Inches(7.6), Inches(4.4))
card(s, Inches(8.5), y0+Inches(0.1), Inches(4.15), Inches(4.1),
     title="Why the 3–4× jump", accent=GOLD, tsize=15,
     bullets=[("", "Does NOT mean per-attribute models learn better"),
              ("", "Per-attribute evaluation runs on a narrower distribution"),
              ("~80–140", " vocabulary words vs ~600+"),
              ("Shared scaffolding:", " ~5 words of constant 4-gram"),
              ("→", " they measure performance on different distributions")], bsize=12)
notes(s, """
A slide of methodological caution, because it's easy to draw the wrong conclusion.
You'll have noticed that the per-attribute numbers — BLEU-4 from 0.33 to 0.47 — are much higher than
the global BLEU-4, around 0.13. A three- or four-fold jump. The temptation is to say: "so the
per-attribute models are much better". Wrong.
That jump doesn't reflect better learning, but the fact that per-attribute evaluation runs on a much
narrower distribution. When you focus on a single attribute, the vocabulary collapses from over 600
words to about 80-140, and the shared scaffolding — "the cheese has a such-and-such attribute of..."
— alone contributes five words of constant 4-gram that send BLEU soaring.
They are two different rulers measuring performance on different distributions. The global model is
intrinsically harder because it must also choose which dimension to describe. Comparing the two
numbers as if on the same scale would be a mistake — and it's the kind of subtlety that separates an
honest analysis from one that fools itself.
""")

# =====================================================================
# SLIDE 25 — SHUFFLE TEST (figure) KEY
# =====================================================================
s = newslide()
chrome(s, 25)
y0 = kicker_title(s, "The key result", "The shuffle test: does the model use the image?")
pic_fit(s, A("fig_shuffle.png"), Inches(0.7), y0, Inches(5.6), Inches(4.4))
card(s, Inches(6.55), y0, Inches(6.1), Inches(2.05),
     title="How it works", accent=PRIM, tsize=15.5,
     body="Predictions are shuffled across test rows, breaking the prediction↔image alignment. "
          "Recompute 100 times → null distribution → z-score. z > 3 ⟺ p < 0.001: strong evidence "
          "of image-conditioning.", bsize=13)
card(s, Inches(6.55), y0+Inches(2.2), Inches(6.1), Inches(2.4),
     title="What it says", accent=RED, tsize=15.5, fill=DEEP, tcolor=GOLD,
     bullets=[("m1 (ResNet):", " z≈0 everywhere → pure language model"),
              ("m3 / m6 (ViT):", " use the image on 4/7 attributes"),
              ("Smell, Thickness, Colour, Structure", " → succeed"),
              ("Aroma", " → no model uses the image")], bsize=12.5)
notes(s, """
We're at the most important result of the whole project. The question is as simple as it is
fundamental: is the model actually using the image, or just mimicking the caption distribution?
The tool is the shuffle test. The idea: if a model is truly conditioned on the image, its prediction
for image i must match the reference of i better than a randomly picked prediction. So we shuffle the
predictions across test rows, breaking the alignment, recompute the overlap a hundred times to get a
null distribution, and measure the z-score. A z above 3 corresponds to p below 0.001: strong
evidence.
The verdict, in the chart, is crystal clear. m1, the ResNet model, has z near zero everywhere: it
never uses the image, it's effectively a pure language model guessing from the distribution. m3 and
m6, the ViT models, clearly use the image on four attributes out of seven: smell, thickness, colour
and structure. The only one where no model uses the image is aroma.
Keep this chart in mind: the next slide draws the architectural conclusion.
""")

# =====================================================================
# SLIDE 26 — ARCHITECTURAL CONCLUSION
# =====================================================================
s = newslide()
chrome(s, 26)
y0 = kicker_title(s, "The interpretation", "The encoder matters, not the decoder")
card(s, Inches(0.7), y0+Inches(0.1), Inches(5.9), Inches(3.5),
     title="m1 vs m3 — isolates the encoder", accent=M3C, tsize=16,
     body="m1 (ResNet) NEVER conditions on the image.\nm3 (ViT) does on most attributes.\n\n"
          "→ A frozen ResNet-50's features aren't enough at this data scale; a frozen ViT-B/16's "
          "are.", bsize=14)
card(s, Inches(6.75), y0+Inches(0.1), Inches(5.9), Inches(3.5),
     title="m3 vs m6 — isolates the decoder", accent=M6C, tsize=16,
     body="They produce different outputs (never the same caption), but have the SAME "
          "image-conditioning profile: the same attributes succeed, the same fail.\n\n"
          "→ The decoder changes the style, not whether the image is used.", bsize=14)
rect(s, Inches(0.7), y0+Inches(3.8), Inches(11.95), Inches(0.75), DEEP, round_=True)
_, tf=textbox(s, Inches(1.0), y0+Inches(3.8), Inches(11.5), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]
setrun(p,"The «binary gate» of image-conditioning is the visual encoder — not the decoder, not data scale, not hyperparameters.",13.5,GOLD,bold=True)
notes(s, """
Put the shuffle test together with the two-axis design from the architectures slide, and the
conclusion writes itself.
Comparing m1 with m3, which isolates the encoder: m1 with ResNet never conditions on the image, m3
with ViT does on most attributes. By changing only the encoder, image use switches on. The reading
is that a frozen ResNet-50's features aren't informative enough at this data scale, while a frozen
ViT-B/16's are.
Comparing m3 with m6, which isolates the decoder: they produce different outputs — they never give
the same caption for the same image — yet they have exactly the same image-conditioning profile. The
same attributes succeed, the same fail. Changing the decoder changes the sentence's style, not
whether the image is used.
The synthesis is at the bottom: the binary gate of image-conditioning, on this dataset, is the
visual encoder. Not the decoder, not the amount of data, not the hyperparameters. It's a clean
result, and it also tells us where to act in future: on the encoder. Now let's reinforce this
conclusion with the complementary metrics.
""")

# =====================================================================
# SLIDE 27 — BEYOND BLEU: 7 METRICS
# =====================================================================
s = newslide()
chrome(s, 27)
y0 = kicker_title(s, "Honest evaluation", "Beyond BLEU: seven complementary metrics")
rows=[("Metric","What it measures","Family"),
      ("BLEU-1/4","n-gram precision (words / exact phrasing)","text"),
      ("METEOR","overlap with stemming and synonyms","text"),
      ("ROUGE-L","longest common subsequence (order)","text"),
      ("CIDEr","TF-IDF weighted n-grams: rewards rare salient terms","text"),
      ("BERTScore","semantic similarity via embeddings (BERT it)","semantic"),
      ("Vocab. conformity","% words in the attested sensory lexicon","domain"),
      ("CLIPScore","caption–image appropriateness (CLIP cosine)","image")]
tx=Inches(0.7); ty=y0; tw=Inches(11.95); th=Inches(4.3)
tbl=s.shapes.add_table(len(rows),3,tx,ty,tw,th).table
for ci,w in enumerate([2.6,7.15,2.2]): tbl.columns[ci].width=Inches(w)
for ri,row in enumerate(rows):
    last = (ri==len(rows)-1)
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ci==2: p.alignment=PP_ALIGN.CENTER
        if ri==0:
            setrun(p,val,13,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            bold = last or ci==0
            setrun(p,val,12.5, (GOLD_D if last else INK), bold=bold)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0xEC,0xE3,0xCF) if last else (CREAM2 if ri%2 else WHITE))
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
notes(s, """
We've seen that BLEU, on its own, misleads: it rewards template repetition and never looks at the
image. To read the models more honestly, we paired BLEU with six complementary metrics, groupable
into three families.
The text metrics compare the prediction with the textual reference: besides BLEU, METEOR — which is
more lenient because it aligns with stemming and synonyms — ROUGE-L, which rewards word order, and
CIDEr, designed specifically for captioning, which weights n-grams with TF-IDF and so rewards rare,
informative sensory terms like "eyes" (occhiatura) over common words.
Then a semantic metric, BERTScore, measuring meaning similarity via embeddings. A domain metric,
vocabulary conformity, which asks: does the model "speak cheese" in the certified register?
And finally, highlighted at the bottom, the most important for captioning: CLIPScore, the only one
that truly looks at the image. It compares the caption with the slice in CLIP space, ignoring the
panel reference. The next two slides show what CIDEr and CLIPScore add.
""")

# =====================================================================
# SLIDE 28 — CIDEr (native chart)
# =====================================================================
s = newslide()
chrome(s, 28)
y0 = kicker_title(s, "CIDEr", "Discriminates where BLEU flattens")
cd = CategoryChartData()
cd.categories = ["random","freq_weighted","m1","m3","m6","most_frequent"]
cd.add_series("Mean CIDEr (7 attributes)", (0.234,0.370,0.586,0.646,0.704,0.790))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), y0, Inches(7.7), Inches(4.35), cd)
ch=gf.chart; style_chart(ch, legend=False)
plot=ch.plots[0]; plot.gap_width=70; plot.has_data_labels=True
plot.data_labels.number_format='0.00'; plot.data_labels.number_format_is_linked=False
plot.data_labels.font.size=Pt(12); plot.data_labels.font.bold=True; plot.data_labels.position=XL_LABEL_POSITION.OUTSIDE_END
points_color(plot.series[0],[GREY,GREY,M1C,M3C,M6C,GOLD])
ch.value_axis.has_major_gridlines=True
card(s, Inches(8.5), y0+Inches(0.1), Inches(4.15), Inches(4.1),
     title="Reading", accent=M6C, tsize=15,
     bullets=[("Orders the models:", " random < freq < m1 < m3 < m6"),
              ("", "Trained models clearly beat the random baselines"),
              ("m6", " is the best among the models"),
              ("most_frequent (0.79)", " leads only because the reference IS often the frequent sentence: artifact, not quality")], bsize=11.5)
notes(s, """
Let's start with CIDEr, the metric designed specifically for captioning. Unlike BLEU, which
flattened the models onto close values, CIDEr weights rare terms with TF-IDF and so separates them.
Look at the ordering in the chart, on the seven-attribute means: random 0.23, then freq_weighted
0.37, then the trained models rising clearly — m1 at 0.59, m3 at 0.65, m6 at 0.70. CIDEr orders the
models correctly, and confirms m6 is the best. It's exactly the information BLEU couldn't give us: a
clean ranking.
There's one column taller than all, in gold: most_frequent, at 0.79. Don't be fooled: the constant
baseline leads only because, on this dataset, the panel reference is often the frequent sentence
itself. It's an artifact of the single reference, not real quality. On an extreme case like rind
thickness, where BLEU gave 0.39 for everyone, CIDEr ranges from 1.04 to 2.29: a rich ordering BLEU
simply did not see.
Conclusion: CIDEr adds real information and separates the trained models from the noise.
""")

# =====================================================================
# SLIDE 29 — CLIPSCORE (double native chart)
# =====================================================================
s = newslide()
chrome(s, 29)
y0 = kicker_title(s, "CLIPScore", "Independent confirmation of the shuffle test")
cd1 = CategoryChartData()
cd1.categories = ["most_freq","random","freq_w","m1","m3","m6"]
cd1.add_series("Mean CLIPScore", (0.1902,0.1913,0.1919,0.1925,0.1921,0.1933))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), y0+Inches(0.25), Inches(5.9), Inches(3.7), cd1)
ch=gf.chart; style_chart(ch, legend=False)
ch.value_axis.minimum_scale=0.185; ch.value_axis.maximum_scale=0.195
plot=ch.plots[0]; plot.gap_width=60
points_color(plot.series[0],[GOLD,GREY,GREY,M1C,M3C,M6C])
_, tf=textbox(s, Inches(0.7), y0-Inches(0.05), Inches(5.9), Inches(0.32))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"By MODEL → all equal (gap < 0.004)",12,DEEP,bold=True)
cd2 = CategoryChartData()
cd2.categories = ["Aroma","Taste","Thickness","Smell","Structure","Colour","Texture"]
cd2.add_series("CLIPScore", (0.181,0.184,0.184,0.185,0.197,0.206,0.208))
gf2=s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.85), y0+Inches(0.25), Inches(5.8), Inches(3.7), cd2)
ch2=gf2.chart; style_chart(ch2, legend=False)
plot2=ch2.plots[0]; plot2.gap_width=50
points_color(plot2.series[0],[PRIM,RED,GOLD,PRIM,GREEN,GREEN,GREEN])
ch2.value_axis.minimum_scale=0.17
_, tf=textbox(s, Inches(6.85), y0-Inches(0.05), Inches(5.8), Inches(0.32))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"By ATTRIBUTE → visual high, smell/taste low",12,DEEP,bold=True)
rect(s, Inches(0.7), y0+Inches(4.05), Inches(11.95), Inches(0.6), DEEP, round_=True)
_, tf=textbox(s, Inches(1.0), y0+Inches(4.05), Inches(11.5), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]
setrun(p,"With a frozen encoder the image isn't exploited: the models do language modeling. The real variation is per attribute.",12.5,GOLD,bold=True)
notes(s, """
Now CLIPScore, the only metric that looks at the image, and the independent confirmation of the
shuffle test. Two charts, two messages.
On the left, mean CLIPScore by model. Note the axis: I deliberately zoomed in, and despite the zoom
the bars are practically identical. The trained models m1, m3, m6 do not beat the constant baseline
most_frequent — the one that always fires the same sentence ignoring the image. The gap is below
0.004, within the noise. It's the quantitative, independent confirmation of what the shuffle test
already told us: with a frozen encoder, the models don't anchor the caption to the visual content,
they do language modeling on the caption distribution.
On the right, the real variation: it's per attribute, not per model. CLIP "sees" that texture,
colour and structure — the visual attributes — get higher scores, while aroma, taste and smell stay
low for everyone, because no image can reveal a taste. It's a sanity check that validates the metric.
And remember the exception: rind thickness is visible but low, because the rind is a small portion
of the slice and CLIP weights it little. Everything fits.
""")

# =====================================================================
# SLIDE 30 — BLEU TRAP (figure)
# =====================================================================
s = newslide()
chrome(s, 30)
y0 = kicker_title(s, "The methodological moral", "The BLEU trap, visualised")
pic_fit(s, A("m_bleu_trap.png"), Inches(0.7), y0, Inches(8.2), Inches(4.4))
card(s, Inches(9.1), y0+Inches(0.1), Inches(3.55), Inches(4.1),
     title="A metric lies", accent=RED, tsize=15.5, fill=DEEP, tcolor=GOLD,
     body="Left, BLEU-1: the constant caption most_frequent WINS.\n\n"
          "Right, CLIPScore: the two models are practically equal.\n\n"
          "→ BLEU-1 alone is misleading. Metrics must be read together.", bsize=13.5)
notes(s, """
This slide is the methodological moral of the whole evaluation, condensed into one image.
Look at the same pair of models — the constant baseline most_frequent against m6 — seen by two
different metrics. On the left, with BLEU-1, the constant caption wins: it always emits the same
sentence and matches on the prefix. On the right, with CLIPScore, the two are practically
indistinguishable.
The exact same comparison, two opposite verdicts. This vividly demonstrates that a single metric, on
its own, can lie. Had we looked only at BLEU-1, we'd have concluded that a constant sentence beats a
trained model — an absurd conclusion. This is exactly why we built a battery of metrics and read
them together: BLEU for phrasing, CIDEr for salient terms, BERTScore for meaning, and CLIPScore for
image grounding.
BERTScore, for completeness, is poorly discriminating here: all values fall between 0.83 and 0.92,
because all sentences talk about cheese with the same structure. It's useful as a sanity check, not
for ranking. The lesson stands: never trust a single metric.
""")

# =====================================================================
# SLIDE 31 — QUALITATIVE EXAMPLES
# =====================================================================
s = newslide()
chrome(s, 31)
y0 = kicker_title(s, "Qualitative examples", "What m6 actually generates")
rows=[("Attribute","Prediction (m6)","Reference (panel)"),
      ("Aroma","The cheese has an aroma of cream.","The cheese has an aroma of cream.  ✓"),
      ("Taste","…taste slightly salty and spicy.","The cheese has a salty taste."),
      ("Thickness","The rind is medium-thick.","The rind is medium-thick.  ✓"),
      ("Colour","…deep, even yellow colour.","…overly deep yellow colour."),
      ("Structure","…irregular fracture and fine grain.","…stretched, with sparse grain in places.")]
tx=Inches(0.7); ty=y0; tw=Inches(11.95); th=Inches(3.3)
tbl=s.shapes.add_table(len(rows),3,tx,ty,tw,th).table
for ci,w in enumerate([1.9,5.0,5.05]): tbl.columns[ci].width=Inches(w)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=""
        p=cell.text_frame.paragraphs[0]
        if ri==0:
            setrun(p,val,12.5,WHITE,bold=True); cell.fill.solid(); cell.fill.fore_color.rgb=DEEP
        else:
            match = "✓" in val
            setrun(p,val,12, (GREEN if match else INK), bold=(ci==0 or match))
            cell.fill.solid(); cell.fill.fore_color.rgb = CREAM2 if ri%2 else WHITE
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
card(s, Inches(0.7), ty+Inches(3.45), Inches(11.95), Inches(0.9),
     title=None, accent=GOLD,
     body="Four recurring patterns: (1) the scaffolding is always correct and worth ~50% of BLEU-4; "
          "(2) the language is fluent thanks to the pretrained LM; (3) the descriptor is often valid "
          "but for another cheese; (4) references vary enormously across tasters → a ceiling on BLEU.",
     bsize=12.5)
notes(s, """
The numbers tell half the story; the generated captions tell the other half. Here you see real
predictions from m6 alongside the panel reference, on test-set samples.
Sometimes it's an exact match: on the aroma of cream, prediction and reference coincide word for
word. On thickness, "the rind is medium-thick", again perfect. But look at the other cases: on taste
the model says "slightly salty and spicy" where the panel said just "salty"; on structure it
proposes "irregular fracture and fine grain" while the reference talks about stretched paste.
From these examples four patterns emerge, summarised at the bottom. First: the scaffolding, "the
cheese has a...", is always correct and alone worth about half the BLEU-4 — which is why the numbers
aren't tiny. Second: the language is fluent and grammatical, thanks to the pretrained model. Third,
and most telling: the descriptor is often plausible but refers to another cheese — exactly what we
expect from a model only partially conditioned on the image. Fourth: the variability across tasters
is enormous and puts a ceiling on what any model can achieve with a single-reference BLEU.
""")

# =====================================================================
# SLIDE 32 — CONCLUSIONS & FUTURE
# =====================================================================
s = newslide(DEEP)
rect(s, 0, 0, Inches(0.22), SH, GOLD)
cheese_holes(s, Inches(11.6), Inches(1.1), GOLD, base=0.55, alpha=16)
_, tf=textbox(s, Inches(0.85), Inches(0.6), Inches(11), Inches(0.9))
setrun(tf.paragraphs[0],"Conclusions & future work",34,CREAM,bold=True,font=FONT_H)
rect(s, Inches(0.88), Inches(1.5), Inches(2.0), Pt(4), GOLD)
def dcard(x,w,title,acc,bullets):
    rect(s, x, Inches(1.85), w, Inches(3.4), RGBColor(0x24,0x40,0x59), round_=True, shadow=True)
    rect(s, x, Inches(1.85), w, Inches(0.1), acc)
    _, tf=textbox(s, x+Inches(0.25), Inches(2.0), w-Inches(0.45), Inches(3.15))
    p=tf.paragraphs[0]; setrun(p,title,15.5,acc,bold=True,font=FONT_H)
    for b in bullets:
        p=tf.add_paragraph(); p.space_before=Pt(6)
        setrun(p,"● ",10,acc,bold=True)
        if isinstance(b,tuple): setrun(p,b[0],12.5,CREAM,bold=True); setrun(p,b[1],12.5,RGBColor(0xD9,0xD2,0xC2))
        else: setrun(p,b,12.5,RGBColor(0xD9,0xD2,0xC2))
dcard(Inches(0.7),Inches(3.85),"Scientific result",GREEN,
      [("The encoder matters:", " ResNet ignores the image, ViT uses it on 4/7"),
       ("", "Independent of decoder, data scale, hyperparameters"),
       ("", "Confirmed by shuffle test + CLIPScore")])
dcard(Inches(4.74),Inches(3.85),"Practical result",GOLD,
      [("4/7", " attributes beat the baseline (max Structure +0.128)"),
       ("m6 (ViT+GePpeTto)", " wins narrowly almost everywhere"),
       ("$5.60", " data cost · 38,437 pairs delivered")])
dcard(Inches(8.78),Inches(3.87),"Future work",PRIM,
      [("Fine-tune the encoder", " and re-measure CLIPScore"),
       ("", "k-NN retrieval baseline as an image-aware «floor»"),
       ("", "If models pull away from the baseline → image exploited")])
_, tf=textbox(s, Inches(0.7), Inches(5.45), Inches(11.9), Inches(0.95))
p=tf.paragraphs[0]
setrun(p,"Honest framing:  ",13,GOLD,bold=True)
setrun(p,"no model would replace a taster. They learned the vocabulary and a partial image→descriptor map — consistent with the data scale and the noise of multi-taster annotation.",13,CREAM)
_, tf=textbox(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.6))
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p,"Thank you. — Questions?",18,GOLD,bold=True,font=FONT_H)
notes(s, """
Let's wrap up. The project delivers two results and one promise for the future.
The scientific result is the most solid: on this dataset, the determinant of image use is the visual
encoder. ResNet ignores the image on all seven attributes, ViT uses it on four out of seven,
independent of the decoder, the data scale and the hyperparameters. And it's not an isolated claim:
two independent tools — the shuffle test and CLIPScore — converge on the same conclusion.
The practical result: the trained models beat the constant baseline on four attributes out of seven,
with the maximum on structure. m6, the ViT with GePpeTto, wins narrowly almost everywhere. And all
this rests on a 38-thousand-pair dataset built with just 5 dollars and 60 cents of LLM cost.
Future work follows directly from the diagnosis: since the bottleneck is the frozen encoder, the
natural next experiment is fine-tuning the encoder and re-measuring CLIPScore. If the models start
pulling away from the constant baseline, we'll have direct proof that the image is finally being
exploited.
I'll close with an honest framing: none of these models would replace a taster today. They learned
the sensory vocabulary and a partial map from image to descriptor — a result consistent with the
data scale and the intrinsic noise of multi-taster annotation. But the path is mapped, and we know
exactly where to act.
Thank you for your attention. I'm happy to take questions.
""")

prs.save(OUTFILE)
print("Saved:", OUTFILE)
print("Total slides:", len(prs.slides._sldIdLst))
